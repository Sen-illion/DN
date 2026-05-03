from __future__ import annotations

import csv
import json
import math
import textwrap
import warnings
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import sys

ROOT = Path(__file__).resolve().parents[1]
VENDOR = ROOT / "vendor_py"
if str(VENDOR) not in sys.path:
    sys.path.insert(0, str(VENDOR))

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


HANDOFF = Path(r"D:\Projects\DN\experiments\handoff_2026-04-30_main_experiments")
OUTDIR = ROOT
ASSETS = ROOT / "assets"
PREVIEWS = ROOT / "previews"
SRC = ROOT / "src"
PPTX_PATH = ROOT / "DN_full_experiment_handoff_detailed_2026-05-01.pptx"
PDF_PATH = ROOT / "DN_full_experiment_handoff_detailed_2026-05-01.pdf"
NOTES_PATH = ROOT / "deck_notes.md"

TABLES = HANDOFF / "01_main_tables"
REVIEWS = HANDOFF / "02_sample_reviews" / "formal20_quality_review"

SLIDE_W = 13.333
SLIDE_H = 7.5
PREVIEW_W = 1600
PREVIEW_H = 900


PALETTE = {
    "bg": "#F6F1E8",
    "paper": "#FFFDF9",
    "ink": "#1E2A39",
    "muted": "#5F6C7B",
    "teal": "#137C72",
    "teal_dark": "#0C5D56",
    "sand": "#E8DAC3",
    "orange": "#D97841",
    "red": "#B64C4C",
    "gold": "#C19A34",
    "line": "#D7C9B2",
    "navy": "#24384D",
    "green": "#50806B",
}


@dataclass
class SlideSpec:
    section: str
    title: str
    conclusion: str
    template: str
    bullets: list[str] | None = None
    table_rows: list[list[str]] | None = None
    table_headers: list[str] | None = None
    asset: str | None = None
    assets: list[str] | None = None
    notes: list[str] | None = None
    small_note: str | None = None


def ensure_dirs() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    PREVIEWS.mkdir(parents=True, exist_ok=True)


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = []
    if bold:
        candidates.extend(
            [
                r"C:\Windows\Fonts\msyhbd.ttc",
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\arialbd.ttf",
            ]
        )
    else:
        candidates.extend(
            [
                r"C:\Windows\Fonts\msyh.ttc",
                r"C:\Windows\Fonts\arial.ttf",
            ]
        )
    for path in candidates:
        p = Path(path)
        if p.exists():
            return ImageFont.truetype(str(p), size=size)
    return ImageFont.load_default()


FONT_REG = font(26)
FONT_BOLD = font(28, bold=True)
FONT_H1 = font(56, bold=True)
FONT_H2 = font(38, bold=True)
FONT_H3 = font(30, bold=True)
FONT_SMALL = font(20)
FONT_TINY = font(16)


def draw_wrapped_text(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, fill: str, fnt, line_gap: int = 8, align: str = "left") -> int:
    x1, y1, x2, y2 = box
    width = x2 - x1
    avg = max(int(getattr(fnt, "size", 20) * 1.1), 20)
    wrap_chars = max(8, int(width / avg * 2))
    lines = []
    for raw in text.split("\n"):
        lines.extend(textwrap.wrap(raw, width=wrap_chars) or [""])
    y = y1
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=fnt)
        line_w = bbox[2] - bbox[0]
        line_h = bbox[3] - bbox[1]
        if align == "center":
            tx = x1 + (width - line_w) / 2
        else:
            tx = x1
        draw.text((tx, y), line, font=fnt, fill=fill)
        y += line_h + line_gap
        if y > y2:
            break
    return y


def add_textbox(slide, left, top, width, height, text, size=20, bold=False, color=None, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei", fill=None):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        tx.fill.solid()
        tx.fill.fore_color.rgb = hex_rgb(fill)
        tx.line.color.rgb = hex_rgb(fill)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color or PALETTE["ink"])
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tx


def add_bullet_box(slide, left, top, width, height, bullets: Iterable[str], size=20):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = hex_rgb(PALETTE["ink"])
        p.bullet = True
    return tx


def add_header(slide, title: str, conclusion: str, section: str):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_rgb(PALETTE["bg"])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(0.28), Inches(12.55), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = hex_rgb(PALETTE["sand"])
    band.line.color.rgb = hex_rgb(PALETTE["sand"])
    add_textbox(slide, 0.52, 0.33, 1.0, 0.3, "本页结论", size=17, bold=True, color=PALETTE["orange"])
    add_textbox(slide, 1.52, 0.33, 10.8, 0.3, conclusion, size=18, color=PALETTE["ink"])
    add_textbox(slide, 0.55, 0.88, 10.5, 0.55, title, size=24, bold=True, color=PALETTE["navy"])
    add_textbox(slide, 11.2, 0.92, 1.5, 0.25, section, size=12, bold=True, color=PALETTE["muted"], align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_rgb(PALETTE["line"])
    line.line.color.rgb = hex_rgb(PALETTE["line"])


def make_bar_chart(path: Path, title: str, labels: list[str], values: list[float], colors: list[str], xlabel: str = "秒（s）", subtitle: str | None = None):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(9.4, 5.2), dpi=180)
    fig.patch.set_facecolor(PALETTE["paper"])
    ax.set_facecolor(PALETTE["paper"])
    y = list(range(len(labels)))
    bars = ax.barh(y, values, color=colors, height=0.58)
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=12)
    ax.invert_yaxis()
    ax.set_xlabel(xlabel, fontsize=11)
    ax.set_title(title, fontsize=16, fontweight="bold", loc="left", color=PALETTE["navy"])
    if subtitle:
        ax.text(0.0, 1.02, subtitle, transform=ax.transAxes, fontsize=10, color=PALETTE["muted"])
    ax.grid(axis="x", color="#D9D1C2", linestyle="--", linewidth=0.8, alpha=0.9)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#C9BCA4")
    max_v = max(values) if values else 1.0
    ax.set_xlim(0, max_v * 1.22)
    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + max_v * 0.02, bar.get_y() + bar.get_height() / 2, f"{val:.3f}".rstrip("0").rstrip("."), va="center", fontsize=11, color=PALETTE["ink"])
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def make_dual_panel_chart(path: Path):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    a_labels = ["DN", "LIGHT", "PWR", "GenAgents"]
    a_first = [7.662, 0.41, 0.833, 9.76]
    a_next = [None, 0.431, 0.837, 6.285]
    b_labels = ["StoryDiffusion", "SDM-v2", "IC-LoRA"]
    b_first = [7.69, 2.847, 30.06]
    b_next = [8.495, 2.827, 30.057]

    fig, axes = plt.subplots(1, 2, figsize=(13.8, 5.8), dpi=180)
    fig.patch.set_facecolor(PALETTE["paper"])

    for ax in axes:
        ax.set_facecolor(PALETTE["paper"])
        ax.spines[["top", "right", "left"]].set_visible(False)
        ax.spines["bottom"].set_color("#C9BCA4")
        ax.grid(axis="x", color="#D9D1C2", linestyle="--", linewidth=0.8, alpha=0.9)

    y1 = list(range(len(a_labels)))
    axes[0].barh([v - 0.18 for v in y1], a_first, height=0.32, color=PALETTE["teal"], label="首个可用响应")
    axes[0].barh([v + 0.18 for v in y1], [v or 0 for v in a_next], height=0.32, color=PALETTE["orange"], label="下一轮响应")
    axes[0].set_yticks(y1)
    axes[0].set_yticklabels(a_labels, fontsize=11)
    axes[0].invert_yaxis()
    axes[0].set_title("A组：文本 / 可玩响应", fontsize=15, fontweight="bold", loc="left", color=PALETTE["navy"])
    axes[0].legend(loc="lower right", frameon=False, fontsize=9)
    axes[0].set_xlim(0, 11.8)
    for idx, (f, n) in enumerate(zip(a_first, a_next)):
        axes[0].text(f + 0.18, idx - 0.18, f"{f:.3f}".rstrip("0").rstrip("."), va="center", fontsize=9)
        if n is not None:
            axes[0].text(n + 0.18, idx + 0.18, f"{n:.3f}".rstrip("0").rstrip("."), va="center", fontsize=9)

    y2 = list(range(len(b_labels)))
    axes[1].barh([v - 0.18 for v in y2], b_first, height=0.32, color=PALETTE["teal"], label="first-turn")
    axes[1].barh([v + 0.18 for v in y2], b_next, height=0.32, color=PALETTE["orange"], label="next-turn")
    axes[1].set_yticks(y2)
    axes[1].set_yticklabels(b_labels, fontsize=11)
    axes[1].invert_yaxis()
    axes[1].set_title("B组：图像剧情延续响应", fontsize=15, fontweight="bold", loc="left", color=PALETTE["navy"])
    axes[1].legend(loc="lower right", frameon=False, fontsize=9)
    axes[1].set_xlim(0, 33.5)
    for idx, (f, n) in enumerate(zip(b_first, b_next)):
        axes[1].text(f + 0.45, idx - 0.18, f"{f:.3f}".rstrip("0").rstrip("."), va="center", fontsize=9)
        axes[1].text(n + 0.45, idx + 0.18, f"{n:.3f}".rstrip("0").rstrip("."), va="center", fontsize=9)

    fig.suptitle("统一主实验总览：DN 作为统一锚点，比较“继续推进故事要等多久”", fontsize=18, fontweight="bold", x=0.02, ha="left", color=PALETTE["navy"])
    plt.tight_layout(rect=(0, 0, 1, 0.94))
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def make_internal_benchmark_chart(path: Path):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    labels = [
        "世界观生成\n默认 20 条均值",
        "世界观生成\n去 council 均值",
        "完整链路\nworldview 中位",
        "完整链路\n选项生成中位",
        "完整链路\n主角生成中位",
    ]
    values = [33.657, 20.031, 9.003, 0.022, 56.546]
    colors = [PALETTE["navy"], PALETTE["teal"], PALETTE["gold"], PALETTE["orange"], PALETTE["red"]]
    fig, ax = plt.subplots(figsize=(10.6, 5.4), dpi=180)
    fig.patch.set_facecolor(PALETTE["paper"])
    ax.set_facecolor(PALETTE["paper"])
    x = range(len(labels))
    bars = ax.bar(x, values, color=colors, width=0.62)
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylabel("秒（s）", fontsize=11)
    ax.set_title("DN 内部链路耗时拆解：真正慢的是前置世界观与主角生成", fontsize=17, fontweight="bold", loc="left", color=PALETTE["navy"])
    ax.grid(axis="y", color="#D9D1C2", linestyle="--", linewidth=0.8)
    ax.spines[["top", "right"]].set_visible(False)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1.2, f"{val:.3f}".rstrip("0").rstrip("."), ha="center", fontsize=10)
    ax.text(0.01, 0.97, "补充稳定性：full_success_rate=1.0，has_image_rate=0.95，option_count_ge_2_rate=1.0", transform=ax.transAxes, va="top", fontsize=10, color=PALETTE["muted"])
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def make_simple_diagram(path: Path, kind: str):
    img = Image.new("RGB", (1600, 900), PALETTE["paper"])
    draw = ImageDraw.Draw(img)

    def box(xy, fill, outline=None, radius=26):
        draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline or fill, width=3)

    if kind == "cover_flow":
        draw.rectangle((0, 0, 1600, 900), fill=PALETTE["navy"])
        draw.ellipse((1060, 100, 1500, 540), fill=PALETTE["teal"])
        draw.ellipse((1180, 360, 1570, 780), fill=PALETTE["orange"])
        draw.rectangle((0, 760, 1600, 900), fill="#1C2836")
        box((92, 180, 580, 330), "#24384D", outline="#5D7389")
        draw_wrapped_text(draw, (125, 212, 542, 305), "主题 -> 世界 -> 选择 -> 下一步", "#FFFFFF", font(34, True), align="center")
        draw_wrapped_text(draw, (92, 355, 730, 430), "一条完整的交互叙事链路，而不是一次性生成一段文本。", "#E6E9EE", font(24))
        steps = [
            (110, 540, 360, 680, "主题输入"),
            (430, 540, 760, 680, "生成世界与角色"),
            (830, 540, 1130, 680, "玩家做选择"),
            (1180, 540, 1470, 680, "系统继续剧情/图像"),
        ]
        for x1, y1, x2, y2, label in steps:
            box((x1, y1, x2, y2), "#F4F0E8")
            draw_wrapped_text(draw, (x1 + 20, y1 + 35, x2 - 20, y2 - 20), label, PALETTE["navy"], font(28, True), align="center")
        for x in [385, 785, 1145]:
            draw.line((x, 610, x + 35, 610), fill="#FFFFFF", width=8)
            draw.polygon([(x + 35, 610), (x + 16, 596), (x + 16, 624)], fill="#FFFFFF")

    elif kind == "experiment_landscape":
        draw_wrapped_text(draw, (40, 40, 1560, 120), "DN 实验全景：不是一组实验，而是四条证据线并行推进", PALETTE["navy"], font(42, True))
        areas = [
            ((70, 180, 760, 410), PALETTE["teal"], "DN 自身历史效率实验", "回答：DN 在文本式 playable latency 协议下和外部方法的相对位置。"),
            ((840, 180, 1530, 410), PALETTE["orange"], "图像 baseline 正式主实验", "回答：在 formal20 与统一 next-turn 协议下，三条图像 baseline 的速度与连续性。"),
            ((70, 470, 760, 720), PALETTE["gold"], "DOC baseline 接入实验", "回答：即使 upstream 不能完整复现，是否仍能产出可进入 DN 比较流程的数据集。"),
            ((840, 470, 1530, 720), PALETTE["red"], "样例审图 / 状态整理 / 汇报整理", "回答：哪些结果能正式汇报，哪些只是 smoke、fallback 或 blocker。"),
        ]
        for rect, fill, title, desc in areas:
            box(rect, fill)
            draw_wrapped_text(draw, (rect[0] + 28, rect[1] + 24, rect[2] - 24, rect[1] + 82), title, "#FFFFFF", font(30, True))
            draw_wrapped_text(draw, (rect[0] + 28, rect[1] + 98, rect[2] - 24, rect[3] - 24), desc, "#FFFFFF", font(22))

    elif kind == "result_tier":
        draw_wrapped_text(draw, (70, 40, 1530, 120), "当前结果层级：正式结果、辅助证据、历史探索，不应混在一起", PALETTE["navy"], font(42, True))
        levels = [
            ((160, 170, 1440, 320), PALETTE["teal"], "正式主结果", "formal20 图像双主表 + formal20 审图包 + DOC formal20 fallback artifact"),
            ((240, 360, 1360, 510), PALETTE["gold"], "可引用的辅助证据", "DN 历史 playable-latency 主表、benchmark 内部效率指标、formal8 过程证据"),
            ((330, 550, 1270, 700), PALETTE["orange"], "保留但降级", "smoke3、blocker、public fallback、real probe、早期集成脚手架"),
        ]
        for rect, fill, title, desc in levels:
            box(rect, fill)
            draw_wrapped_text(draw, (rect[0] + 28, rect[1] + 24, rect[2] - 24, rect[1] + 82), title, "#FFFFFF", font(32, True), align="center")
            draw_wrapped_text(draw, (rect[0] + 46, rect[1] + 90, rect[2] - 46, rect[3] - 22), desc, "#FFFFFF", font(24), align="center")

    elif kind == "measure_protocol":
        draw_wrapped_text(draw, (60, 40, 1540, 110), "统一主实验在测什么：虽然输出形态不同，但都在测“用户要等多久，故事才能继续”", PALETTE["navy"], font(40, True))
        box((90, 190, 740, 680), "#EAF5F2", outline=PALETTE["teal"])
        box((860, 190, 1510, 680), "#FBEEE5", outline=PALETTE["orange"])
        draw_wrapped_text(draw, (130, 230, 700, 290), "A组：文本 / 可玩响应", PALETTE["teal_dark"], font(34, True))
        draw_wrapped_text(draw, (130, 315, 700, 600), "测量口径：\n1. 玩家收到第一段可继续玩的响应用了多久\n2. 下一轮继续响应用了多久\n\n代表系统：DN、LIGHT、PWR、GenAgents", PALETTE["ink"], font(24))
        draw_wrapped_text(draw, (900, 230, 1470, 290), "B组：图像剧情延续响应", PALETTE["orange"], font(34, True))
        draw_wrapped_text(draw, (900, 315, 1470, 600), "测量口径：\n1. first-turn 首张剧情图生成用了多久\n2. next-turn 玩家点击动作后，下一张剧情图用了多久\n\n代表系统：StoryDiffusion、SDM-v2、IC-LoRA", PALETTE["ink"], font(24))
        draw_wrapped_text(draw, (170, 620, 1430, 760), "共同上位问题：当故事要继续推进时，用户等待系统返回“下一步内容”的时间是多少。", PALETTE["navy"], font(26, True), align="center")

    elif kind == "doc_status":
        draw_wrapped_text(draw, (50, 40, 1540, 110), "DOC baseline 当前状态：已经能进比较链路，但不是 upstream 全栈复现", PALETTE["navy"], font(40, True))
        cols = [
            ((70, 180, 500, 700), PALETTE["teal"], "已完成", ["DN item -> DOC-style premise 映射", "smoke3 / formal8 / formal20 artifact", "统一 result.json / summary.json / manifest"]),
            ((585, 180, 1015, 700), PALETTE["gold"], "结构已搭建", ["run_doc.py adapter runner", "normalized copy 输出路径", "playable schema 对齐与 notes 透明标注"]),
            ((1100, 180, 1530, 700), PALETTE["red"], "尚未完成", ["GPT3/Alpa/OPT upstream 完整复现", "真实官方模型成功跑批", "把 fallback 误写成真实 DOC 输出"]),
        ]
        for rect, fill, title, items in cols:
            box(rect, fill)
            draw_wrapped_text(draw, (rect[0] + 22, rect[1] + 24, rect[2] - 22, rect[1] + 74), title, "#FFFFFF", font(30, True), align="center")
            y = rect[1] + 110
            for item in items:
                draw.ellipse((rect[0] + 30, y + 10, rect[0] + 44, y + 24), fill="#FFFFFF")
                draw_wrapped_text(draw, (rect[0] + 58, y, rect[2] - 24, y + 80), item, "#FFFFFF", font(22))
                y += 110

    elif kind == "takeover_flow":
        draw_wrapped_text(draw, (60, 40, 1540, 110), "结论与后续工作：当前已经形成一条可汇报、可复用、可继续扩展的主线", PALETTE["navy"], font(38, True))
        steps = [
            ("1", "正式主结果已形成", "formal20 图像双主表已稳定"),
            ("2", "文本比较链路已补齐", "DOC fallback + DN 历史 latency"),
            ("3", "系统瓶颈已定位", "世界观与主角生成最重"),
            ("4", "后续优先工作", "扩展正式样本 / 优化端到端时延"),
        ]
        x = 110
        for idx, (num, title, desc) in enumerate(steps):
            box((x, 270, x + 300, 650), "#F8F3EA", outline=PALETTE["line"])
            draw.ellipse((x + 108, 185, x + 192, 269), fill=PALETTE["teal"] if idx < 2 else PALETTE["orange"])
            draw_wrapped_text(draw, (x + 116, 199, x + 184, 255), num, "#FFFFFF", font(34, True), align="center")
            draw_wrapped_text(draw, (x + 26, 315, x + 274, 400), title, PALETTE["navy"], font(26, True), align="center")
            draw_wrapped_text(draw, (x + 26, 430, x + 274, 610), desc, PALETTE["muted"], font(20), align="center")
            if idx < 3:
                draw.line((x + 300, 460, x + 350, 460), fill=PALETTE["navy"], width=8)
                draw.polygon([(x + 350, 460), (x + 326, 444), (x + 326, 476)], fill=PALETTE["navy"])
            x += 350

    elif kind == "appendix_divider":
        draw.rectangle((0, 0, 1600, 900), fill=PALETTE["navy"])
        draw.ellipse((1080, 160, 1480, 560), fill=PALETTE["teal"])
        draw_wrapped_text(draw, (120, 250, 1000, 360), "附录：原始路径、历史索引、补充证据与已知问题", "#FFFFFF", font(50, True))
        draw_wrapped_text(draw, (126, 400, 920, 520), "这一部分用于补充汇报证据和边界条件，不与正式主结果混用。", "#DDE3EA", font(24))

    elif kind == "asset_map":
        draw_wrapped_text(draw, (60, 40, 1540, 120), "实验资产目录图：这套整理目录可以作为汇报材料与原始证据的统一入口", PALETTE["navy"], font(40, True))
        tree = [
            "handoff_2026-04-30_main_experiments/",
            "  00_overview/      -> 一页摘要、结果地图、阅读顺序",
            "  01_main_tables/   -> 正式主表、内部 benchmark 摘要",
            "  02_sample_reviews/-> formal20 审图包与 review index",
            "  03_raw_artifacts/ -> 图像 formal20 / DOC / DN run 原始目录索引",
            "  04_repro_entrypoints/ -> runner、subset、schema、远程入口",
            "  05_status_and_history/ -> ground truth、known issues、日志",
        ]
        box((140, 180, 1460, 720), "#FBF7F0", outline=PALETTE["line"])
        y = 225
        for line in tree:
            draw_wrapped_text(draw, (190, y, 1410, y + 60), line, PALETTE["ink"], font(25, "->" not in line))
            y += 72

    elif kind == "dont_mix":
        draw_wrapped_text(draw, (60, 40, 1540, 110), "哪些结果不能混用：正式结果要保持口径干净", PALETTE["navy"], font(40, True))
        box((110, 190, 720, 700), "#EAF5F2", outline=PALETTE["teal"])
        box((880, 190, 1490, 700), "#FBEEE5", outline=PALETTE["orange"])
        draw_wrapped_text(draw, (170, 230, 660, 290), "可以直接放进主实验叙事", PALETTE["teal_dark"], font(30, True), align="center")
        ok_items = [
            "formal20 first_turn 图像主表",
            "formal20 next_turn 图像主表",
            "formal20 quality review",
            "DOC formal20 fallback artifact（需明确是 fallback）",
            "DN 历史 playable-latency 主表",
        ]
        y = 330
        for item in ok_items:
            draw.ellipse((160, y + 10, 178, y + 28), fill=PALETTE["teal"])
            draw_wrapped_text(draw, (194, y, 645, y + 70), item, PALETTE["ink"], font(23))
            y += 72
        draw_wrapped_text(draw, (950, 230, 1420, 290), "不要与正式主表混在一起", PALETTE["orange"], font(30, True), align="center")
        bad_items = [
            "smoke3 / formal8 过程结果",
            "strict SDM-v2 blocker run",
            "public SD v1.4 sanity fallback",
            "DOC normalized copy 与 raw fallback 的语义差异",
            "image_baseline_summary_20260430 混合表",
        ]
        y = 330
        for item in bad_items:
            draw.line((930, y + 18, 950, y + 38), fill=PALETTE["red"], width=6)
            draw.line((950, y + 18, 930, y + 38), fill=PALETTE["red"], width=6)
            draw_wrapped_text(draw, (972, y, 1415, y + 70), item, PALETTE["ink"], font(23))
            y += 72

    img.save(path)


def build_assets():
    configure_matplotlib()
    make_dual_panel_chart(ASSETS / "unified_overview_panel.png")
    make_bar_chart(
        ASSETS / "a_group_results.png",
        "A组：文本 / 可玩响应速度",
        ["DN", "LIGHT", "PWR", "GenAgents"],
        [7.662, 0.41, 0.833, 9.76],
        [PALETTE["navy"], PALETTE["teal"], PALETTE["gold"], PALETTE["orange"]],
        subtitle="首个可继续游玩的响应时间，DN 为统一锚点。",
    )
    make_bar_chart(
        ASSETS / "b_first_turn.png",
        "B组：first-turn 首张剧情图速度",
        ["StoryDiffusion", "SDM-v2", "IC-LoRA"],
        [7.69, 2.847, 30.06],
        [PALETTE["teal"], PALETTE["gold"], PALETTE["red"]],
        subtitle="formal20，统一 DN-style prompt 适配。",
    )
    make_bar_chart(
        ASSETS / "b_next_turn.png",
        "B组：next-turn 点击后下一张剧情图速度",
        ["StoryDiffusion", "SDM-v2", "IC-LoRA"],
        [8.495, 2.827, 30.057],
        [PALETTE["teal"], PALETTE["gold"], PALETTE["red"]],
        subtitle="formal20，统一动作模板：Choose the most direct action...",
    )
    make_internal_benchmark_chart(ASSETS / "dn_internal_benchmark.png")
    for name in [
        "cover_flow",
        "experiment_landscape",
        "result_tier",
        "measure_protocol",
        "doc_status",
        "takeover_flow",
        "appendix_divider",
        "asset_map",
        "dont_mix",
    ]:
        make_simple_diagram(ASSETS / f"{name}.png", name)


def load_data():
    first_turn = pd.read_csv(TABLES / "first_turn_formal20_summary.csv")
    next_turn = pd.read_csv(TABLES / "next_turn_formal20_latency_summary.csv")
    benchmark = pd.read_csv(TABLES / "benchmark_v1_summary_metrics.csv")
    return first_turn, next_turn, benchmark


def configure_matplotlib():
    import matplotlib
    matplotlib.use("Agg")
    matplotlib.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    matplotlib.rcParams["axes.unicode_minus"] = False
    warnings.filterwarnings("ignore", message="Glyph .* missing from font")


def build_slide_specs(first_turn, next_turn, benchmark) -> list[SlideSpec]:
    unified_rows = [
        ["DN", "A组 文本 / 可玩响应", "7.662s", "-", "20", "DN 自身 first-playable proxy"],
        ["LIGHT", "A组 文本 / 可玩响应", "0.41s", "0.431s", "8", "权威互动文本基线"],
        ["PWR", "A组 文本 / 可玩响应", "0.833s", "0.837s", "5", "速度参考型故事生成基线"],
        ["GenAgents", "A组 文本 / 可玩响应", "9.76s", "6.285s", "8", "连续性 / 多轮状态补充基线"],
        ["StoryDiffusion", "B组 图像剧情延续", "7.69s", "8.495s", "20", "图像连续剧情基线"],
        ["SDM-v2", "B组 图像剧情延续", "2.847s", "2.827s", "20", "单图基础生成基线"],
        ["IC-LoRA", "B组 图像剧情延续", "30.06s", "30.057s", "20", "官方 workflow 较重的连续图基线"],
    ]

    a_rows = [
        ["系统", "首个可玩响应", "下一轮响应", "样本量", "状态"],
        ["DN", "7.662s", "-", "20", "ready"],
        ["LIGHT", "0.41s", "0.431s", "8", "ready"],
        ["PWR", "0.833s", "0.837s", "5", "ready"],
        ["GenAgents", "9.76s", "6.285s", "8", "ready"],
    ]
    b_rows = [
        ["baseline", "first-turn", "next-turn", "sample_size", "success_rate"],
        ["StoryDiffusion", "7.69s", "8.495s", "20", "1.0"],
        ["SDM-v2", "2.847s", "2.827s", "20", "1.0"],
        ["IC-LoRA", "30.06s", "30.057s", "20", "1.0"],
    ]
    history_rows = [
        ["阶段", "代表 run", "状态", "用途"],
        ["smoke3", "StoryDiffusion smoke3 / SDM-v2 blocker / IC-LoRA blocker", "已保留", "验证链路、记录 blocker"],
        ["formal8", "StoryDiffusion formal8 / public SD v1.4 formal8", "已保留", "过渡到 formal20 前的中间证据"],
        ["formal20", "StoryDiffusion / SDM-v2 / IC-LoRA", "正式主结果", "当前图像主实验 ground truth"],
    ]
    doc_rows = [
        ["层级", "当前状态", "能不能直接用"],
        ["DOC fallback artifact", "已完成 smoke3 / formal8 / formal20", "可以进入 DN comparison pipeline"],
        ["DOC normalized copy", "已对齐 unified schema 与 playable schema", "可以用于统一下游读取"],
        ["DOC upstream full stack", "未完整复现 GPT3/Alpa/OPT", "不能当成真实官方复现结果"],
    ]
    benchmark_rows = [
        ["指标", "值"],
        ["worldview_default_20 mean_s", "33.657"],
        ["worldview_no_council_20 mean_s", "20.031"],
        ["fullchain worldview_median_s", "9.003"],
        ["generate_option_median_s", "0.022"],
        ["main_character_median_s", "56.546"],
        ["full_success_rate", "1.0"],
        ["has_image_rate", "0.95"],
        ["option_count_ge_2_rate", "1.0"],
    ]
    formal20_path_rows = [
        ["run", "路径"],
        ["StoryDiffusion first_turn", "remote_baseline_results_20260430/outputs/storydiffusion_formal20/..."],
        ["StoryDiffusion next_turn", "remote_baseline_results_20260430/outputs/storydiffusion_nextturn_formal20/..."],
        ["SDM-v2 first_turn", "remote_baseline_results_20260430/outputs/sdmv2_local_formal20/..."],
        ["SDM-v2 next_turn", "remote_baseline_results_20260430/outputs/sdmv2_nextturn_formal20/..."],
        ["IC-LoRA first_turn", "remote_baseline_results_20260430/outputs/iclora_formal20_real/..."],
        ["IC-LoRA next_turn", "remote_baseline_results_20260430/outputs/iclora_nextturn_formal20_real/..."],
    ]
    repro_rows = [
        ["类别", "入口"],
        ["StoryDiffusion runner", "scripts/baselines/run_storydiffusion.py"],
        ["SDM-v2 runner", "scripts/baselines/run_sdmv2.py"],
        ["IC-LoRA runner", "scripts/baselines/run_iclora.py"],
        ["DOC runner", "scripts/baselines/run_doc.py"],
        ["subset", "baselines/subsets/dn_style_formal20.json"],
        ["schemas", "baseline_integration/schema/*.md"],
        ["远程环境", "docs/AUTODL_CLOUD_EXPERIMENT_ENV.md"],
    ]
    issue_rows = [
        ["问题", "说明"],
        ["StoryDiffusion vs ComfyUI 显存互斥", "next-turn 稳定前需要停掉 ComfyUI，再分阶段重启。"],
        ["SDM-v2 官方仓库下线", "改为本地模型目录加载，而不是直接 online repo id。"],
        ["DOC 不是 upstream 全栈复现", "当前是 faithful fallback adapter，可比较但不能夸大。"],
        ["中文字段乱码", "部分历史材料和元数据有 mojibake，不影响主结果数字。"],
        ["review 图片曾损坏", "已存在 redownload replacement，review 使用替代图。"],
    ]

    slides = [
        SlideSpec("研究汇报", "DN 项目实验全景总结", "这份汇报聚焦三件事：DN 是什么、实验如何设计、当前最正式的结果是什么。", "cover", asset="cover_flow.png", small_note="2026-05-01 · DN 实验汇报版"),
        SlideSpec("研究汇报", "DN 是什么", "DN 的目标不是一次性生成文本，而是持续推进可交互叙事。", "bullets_asset", bullets=["输入可以是一个主题、一段设定，或一个待展开的故事想法。", "系统需要先组织世界、角色与剧情，再返回一个可继续游玩的当前状态。", "玩家做出选择后，系统还要继续生成下一轮文本或图像内容。"], asset="cover_flow.png"),
        SlideSpec("研究汇报", "汇报结构", "主 deck 讲问题、方法、结果与结论；附录只负责补充证据。", "bullets_asset", bullets=["前 18 页先回答：DN 在做什么、实验在比什么、哪些结果最值得引用。", "附录再补充：原始路径、历史 run、schema、远程环境与已知问题。", "正式结果与 smoke / fallback / blocker 已经分层，不能混成同一口径。"], asset="asset_map.png"),
        SlideSpec("研究汇报", "研究问题", "这些实验不是为了堆 baseline，而是为了回答 DN 当前的真实进展。", "bullets_only", bullets=["问题 1：DN 自己的完整链路是否稳定，时间主要花在哪。", "问题 2：如果和外部文本方法比较，DN 的 playable latency 处在什么位置。", "问题 3：如果和图像方法比较，哪些 baseline 已经真正跑通并形成正式对照。"]),
        SlideSpec("研究汇报", "实验全景地图", "目前形成的是四条证据线并行推进，而不是单一目录里的零散结果。", "full_asset", asset="experiment_landscape.png", small_note="四条线分别回答：自身效率、图像对照、DOC 接入、样例与状态整理。"),
        SlideSpec("研究汇报", "当前正式结果边界", "formal20 图像双主表是当前最正式的主结果；其余材料用于补证据与解释边界。", "full_asset", asset="result_tier.png", small_note="正式主结果、辅助证据、历史探索已经分层。"),
        SlideSpec("研究汇报", "主实验到底在比什么", "A组与 B组形态不同，但都在回答“用户要等多久，故事才能继续”。", "full_asset", asset="measure_protocol.png", small_note="重要口径：两组不是完全同构系统，不能被讲成单一竞技榜单。"),
        SlideSpec("研究汇报", "统一主实验总表", "这一页是汇报核心：DN 作为统一锚点，把两类对照组放进同一上位问题。", "table", table_headers=["系统 / baseline", "类别", "首个响应", "下一轮响应", "样本量", "说明"], table_rows=unified_rows, small_note="显著口径说明：A组衡量可继续游玩的文本响应；B组衡量点击后下一张剧情图响应。"),
        SlideSpec("研究汇报", "统一主实验总览图", "整体看，SDM-v2 最快，IC-LoRA 最慢；DN 在文本组中处于中等偏慢位置。", "full_asset", asset="unified_overview_panel.png", small_note="这张图展示的是“故事继续推进时，用户等待时间的量级分布”。"),
        SlideSpec("研究汇报", "A组：文本 / 可玩响应组在测什么", "A组更适合回答“可玩响应效率”，而不是图像连续性。", "table_bullets", bullets=["DN：我们的系统，强调完整交互链路。", "LIGHT：权威互动文本 / game-world 参考基线。", "PWR：速度参考型故事生成基线。", "GenAgents：多轮状态维持与连续性补充基线。"], table_rows=a_rows, table_headers=None, small_note="样本量并不完全相同，因此更适合做相对位置参考。"),
        SlideSpec("研究汇报", "A组结果解读", "外部文本系统里存在更快的方法，但 DN 的价值不只在速度，还在完整系统形态。", "bullets_asset", bullets=["LIGHT 和 PWR 响应更快，但它们和 DN 的语义形态并不完全等同。", "GenAgents 首轮更慢、下一轮更快，更像连续性补充行。", "DN 的 7.662s 应理解为“当前完整链路下的 first-playable proxy”。"], asset="a_group_results.png"),
        SlideSpec("研究汇报", "B组：图像剧情延续组在测什么", "B组聚焦图像剧情延续，最贴近“点击后继续出下一张图”的体验。", "bullets_only", bullets=["统一样本集：formal20。", "统一 next-turn 协议：固定动作模板，避免 baseline 自带交互不一致。", "统一输入风格：DN-style prompt adaptation，保证主题与冲突大体可比。", "统一输出：result.json / summary.json / 样例图。"]),
        SlideSpec("研究汇报", "B组结果：first-turn", "首张剧情图速度上，SDM-v2 最快，StoryDiffusion 居中，IC-LoRA 明显最慢。", "bullets_asset", bullets=["SDM-v2：2.847s，速度优势最明显。", "StoryDiffusion：7.69s，在连续图能力和时延之间较平衡。", "IC-LoRA：30.06s，workflow 更重，因此速度成本最高。"], asset="b_first_turn.png"),
        SlideSpec("研究汇报", "B组结果：next-turn", "next-turn 是最接近真实交互体验的指标：点击动作后，要等下一张图多久。", "bullets_asset", bullets=["StoryDiffusion：8.495s。", "SDM-v2：2.827s。", "IC-LoRA：30.057s。", "排序与 first-turn 基本一致，说明系统开销结构相对稳定。"], asset="b_next_turn.png"),
        SlideSpec("研究汇报", "B组稳定性、连续性与样例图", "formal20 上三条图像 baseline 都已形成稳定、可追溯、可审图的正式结果包。", "image_grid", assets=["DNQBV1_001_review.png", "DNQBV1_010_review.png", "DNQBV1_020_review.png"], bullets=["formal20：3 条 baseline 都达到 20/20 success。", "next_turn continuation_success_rate = 1.0。", "interaction_continuity = 1.0。"], small_note="样例 sheet 的行：first_turn / next_turn current / next_turn next；列：StoryDiffusion / SDM-v2 / IC-LoRA。"),
        SlideSpec("研究汇报", "DOC baseline 当前进展", "DOC 的现实价值在于：已经能稳定产出可比较数据集，而不是 upstream 全栈已复现。", "full_asset", asset="doc_status.png", small_note="正式可用表述：faithful DOC-style fallback artifact for DN comparison pipeline。"),
        SlideSpec("研究汇报", "DN 内部 benchmark 说明了什么", "DN 当前真正慢的主要是世界观与主角生成，而不是“生成选项”这一步。", "bullets_asset", bullets=["世界观默认均值 33.657s；去 council 后仍有 20.031s。", "generate_option 中位仅 0.022s，不是主瓶颈。", "main_character 中位 56.546s，说明部分前置生成很重。"], asset="dn_internal_benchmark.png"),
        SlideSpec("研究汇报", "结论与后续工作", "当前最稳的汇报主线是 formal20 图像主表 + DOC fallback + DN 历史 latency 表。", "full_asset", asset="takeover_flow.png", small_note="如果继续扩展，优先沿用 formal20 协议与现有 runner。"),
        SlideSpec("补充材料", "附录", "附录页用于补充原始路径、历史过程与已知问题。", "full_asset", asset="appendix_divider.png"),
        SlideSpec("补充材料", "实验资产目录图", "这套整理目录就是当前汇报材料与原始证据的统一入口。", "full_asset", asset="asset_map.png"),
        SlideSpec("补充材料", "formal20 图像 baseline raw artifact 路径总表", "正式图像主实验的 source-of-truth 在这些目录，不需要再去猜。", "table", table_headers=formal20_path_rows[0], table_rows=formal20_path_rows[1:], small_note="完整路径见 CURRENT_GROUND_TRUTH.md 与 image_baselines_formal20/README.md。"),
        SlideSpec("补充材料", "smoke3 / formal8 / blocker / fallback 历史图像实验索引", "历史 run 需要保留，但只能作为过程证据，不能替代 formal20 主表。", "table", table_headers=history_rows[0], table_rows=history_rows[1:], small_note="重点是解释为什么最后 formal20 主表会固定成 StoryDiffusion / SDM-v2 / IC-LoRA。"),
        SlideSpec("补充材料", "DOC raw fallback 与 normalized copy 对照", "两个目录都保留：一个强调原始产物，一个强调统一 schema 可读取。", "table", table_headers=doc_rows[0], table_rows=doc_rows[1:], small_note="推荐主讲法：可用，但 fallback-oriented，而非 upstream-authentic。"),
        SlideSpec("补充材料", "DN 历史 playable-latency 原表摘录", "如果要谈 DN 和文本基线的历史对比，这页是最短入口。", "table", table_headers=a_rows[0], table_rows=a_rows[1:], small_note="WorldGeneration 行是 supplementary，不在主表核心 4 行里。"),
        SlideSpec("补充材料", "DN benchmark 内部指标原表摘录", "这页用来支撑“时间到底花在哪”。", "table", table_headers=benchmark_rows[0], table_rows=benchmark_rows[1:], small_note="source-of-truth：benchmark_v1_summary_metrics.csv / .json / .xlsx。"),
        SlideSpec("补充材料", "当前已知问题", "已知问题定义了当前结果的解释边界。", "table", table_headers=issue_rows[0], table_rows=issue_rows[1:], small_note="最重要的边界：SDM-v2 仓库下线、DOC 不是 full upstream reproduction。"),
        SlideSpec("补充材料", "远程 4090 环境与复现入口", "如果继续复跑实验，先找 runner、subset、schema，再进入远程 4090 环境。", "table", table_headers=repro_rows[0], table_rows=repro_rows[1:], small_note="当前远程工作根：/root/autodl-tmp/outputs。"),
        SlideSpec("补充材料", "代表样本说明页", "看样例时，最重要的是横向比 baseline，纵向比 first-turn 与 next-turn 连续性。", "image_grid", assets=["DNQBV1_005_review.png", "DNQBV1_015_review.png"], bullets=["行含义：first_turn、next_turn current、next_turn next。", "列含义：StoryDiffusion、SDM-v2、IC-LoRA。", "适合看 continuity、首图画质、角色/场景延续。"], small_note="已知替代图：SDM-v2 DNQBV1_009、IC-LoRA DNQBV1_007。"),
        SlideSpec("补充材料", "哪些结果不能混用", "把 smoke、blocker、public fallback 与 formal20 正式结果混在一起，会直接破坏口径。", "full_asset", asset="dont_mix.png"),
    ]
    return slides


def add_picture(slide, image_path: Path, left: float, top: float, width: float, height: float):
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_native_table(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[str]], font_size=12):
    row_count = len(rows) + 1
    col_count = len(headers)
    table = slide.shapes.add_table(row_count, col_count, Inches(left), Inches(top), Inches(width), Inches(height)).table
    col_width = width / col_count
    for idx in range(col_count):
        table.columns[idx].width = Inches(col_width)
    header_fill = hex_rgb(PALETTE["navy"])
    for idx, text in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    row_fills = [PALETTE["paper"], "#F2ECE2"]
    for r, row in enumerate(rows, start=1):
        fill = hex_rgb(row_fills[(r - 1) % 2])
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = hex_rgb(PALETTE["ink"])
                p.alignment = PP_ALIGN.CENTER if c != col_count - 1 else PP_ALIGN.LEFT
    return table


def build_ppt(slides: list[SlideSpec]):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        if spec.template == "cover":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = hex_rgb(PALETTE["navy"])
            add_picture(slide, ASSETS / spec.asset, 0, 0, SLIDE_W, SLIDE_H)
            add_textbox(slide, 0.62, 0.55, 5.8, 0.5, "DN 项目实验全景总结", size=28, bold=True, color="#FFFFFF")
            add_textbox(slide, 0.68, 1.15, 6.5, 1.1, "从系统能力、实验设计到正式结果与关键结论", size=17, color="#E2E7EE")
            add_textbox(slide, 0.68, 6.62, 4.5, 0.25, spec.small_note or "", size=11, color="#C8D1DC")
            add_textbox(slide, 10.55, 6.62, 2.15, 0.25, "DN experiment report", size=11, color="#C8D1DC", align=PP_ALIGN.RIGHT)
            continue

        add_header(slide, spec.title, spec.conclusion, spec.section)
        if spec.template == "bullets_only":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.85), Inches(12.0), Inches(4.7))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_rgb(PALETTE["paper"])
            shape.line.color.rgb = hex_rgb(PALETTE["line"])
            add_bullet_box(slide, 1.0, 2.15, 10.8, 3.9, spec.bullets or [], size=22)
        elif spec.template == "bullets_asset":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.82), Inches(5.1), Inches(4.95))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_rgb(PALETTE["paper"])
            shape.line.color.rgb = hex_rgb(PALETTE["line"])
            add_bullet_box(slide, 0.98, 2.1, 4.45, 3.95, spec.bullets or [], size=19)
            add_picture(slide, ASSETS / spec.asset if spec.asset and (ASSETS / spec.asset).exists() else REVIEWS / spec.asset, 6.1, 1.93, 6.15, 4.55)
        elif spec.template == "full_asset":
            add_picture(slide, ASSETS / spec.asset if spec.asset and (ASSETS / spec.asset).exists() else REVIEWS / spec.asset, 0.78, 1.85, 11.9, 4.95)
        elif spec.template == "table":
            add_native_table(slide, 0.72, 1.88, 12.0, 4.95, spec.table_headers or [], spec.table_rows or [], font_size=12)
        elif spec.template == "table_bullets":
            add_native_table(slide, 0.74, 1.95, 5.0, 4.6, spec.table_rows[0], spec.table_rows[1:], font_size=12)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.92), Inches(6.1), Inches(4.65))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_rgb(PALETTE["paper"])
            shape.line.color.rgb = hex_rgb(PALETTE["line"])
            add_bullet_box(slide, 6.25, 2.2, 5.4, 3.8, spec.bullets or [], size=20)
        elif spec.template == "image_grid":
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(1.88), Inches(3.1), Inches(4.88))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_rgb(PALETTE["paper"])
            shape.line.color.rgb = hex_rgb(PALETTE["line"])
            add_bullet_box(slide, 1.02, 2.15, 2.45, 3.8, spec.bullets or [], size=18)
            imgs = spec.assets or []
            positions = [(4.1, 1.9, 4.0, 2.2), (8.25, 1.9, 4.0, 2.2), (4.1, 4.2, 8.15, 2.3)]
            for pos, name in zip(positions, imgs):
                add_picture(slide, REVIEWS / name, *pos)
        if spec.small_note:
            add_textbox(slide, 0.84, 6.72, 11.9, 0.24, spec.small_note, size=11, color=PALETTE["muted"])

    prs.save(PPTX_PATH)


def preview_table(draw: ImageDraw.ImageDraw, box, headers, rows):
    x1, y1, x2, y2 = box
    cols = len(headers)
    rows_n = len(rows) + 1
    col_w = (x2 - x1) / cols
    row_h = (y2 - y1) / rows_n
    draw.rounded_rectangle(box, radius=24, fill=PALETTE["paper"], outline=PALETTE["line"], width=3)
    for c, header in enumerate(headers):
        cx1 = x1 + c * col_w
        cx2 = cx1 + col_w
        draw.rectangle((cx1, y1, cx2, y1 + row_h), fill=PALETTE["navy"])
        draw_wrapped_text(draw, (int(cx1 + 10), int(y1 + 10), int(cx2 - 10), int(y1 + row_h - 10)), str(header), "#FFFFFF", font(18, True), line_gap=4, align="center")
    for r, row in enumerate(rows):
        ry1 = y1 + (r + 1) * row_h
        fill = PALETTE["paper"] if r % 2 == 0 else "#F2ECE2"
        for c, value in enumerate(row):
            cx1 = x1 + c * col_w
            cx2 = cx1 + col_w
            draw.rectangle((cx1, ry1, cx2, ry1 + row_h), fill=fill, outline=PALETTE["line"])
            draw_wrapped_text(draw, (int(cx1 + 8), int(ry1 + 8), int(cx2 - 8), int(ry1 + row_h - 8)), str(value), PALETTE["ink"], font(16), line_gap=4, align="center" if c < cols - 1 else "left")


def build_previews(slides: list[SlideSpec]):
    for idx, spec in enumerate(slides, start=1):
        img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), PALETTE["bg"])
        draw = ImageDraw.Draw(img)
        if spec.template == "cover":
            hero = Image.open(ASSETS / spec.asset).convert("RGB").resize((PREVIEW_W, PREVIEW_H))
            img.paste(hero, (0, 0))
            draw_wrapped_text(draw, (74, 60, 900, 130), "DN 项目实验全景总结", "#FFFFFF", font(42, True))
            draw_wrapped_text(draw, (78, 138, 850, 220), "从系统能力、实验设计到正式结果与关键结论", "#E7ECF2", font(20))
            draw_wrapped_text(draw, (78, 828, 580, 870), spec.small_note or "", "#C9D2DD", font(14))
        else:
            draw.rounded_rectangle((42, 28, 1558, 92), radius=22, fill=PALETTE["sand"])
            draw_wrapped_text(draw, (70, 44, 205, 78), "本页结论", PALETTE["orange"], font(18, True))
            draw_wrapped_text(draw, (210, 44, 1485, 80), spec.conclusion, PALETTE["ink"], font(17))
            draw_wrapped_text(draw, (70, 110, 1240, 160), spec.title, PALETTE["navy"], font(28, True))
            draw_wrapped_text(draw, (1380, 114, 1520, 142), spec.section, PALETTE["muted"], font(12, True), align="center")
            draw.line((70, 172, 1530, 172), fill=PALETTE["line"], width=2)
            if spec.template == "bullets_only":
                draw.rounded_rectangle((80, 220, 1520, 800), radius=26, fill=PALETTE["paper"], outline=PALETTE["line"], width=3)
                y = 270
                for bullet in spec.bullets or []:
                    draw.ellipse((130, y + 8, 145, y + 23), fill=PALETTE["teal"])
                    draw_wrapped_text(draw, (165, y, 1450, y + 85), bullet, PALETTE["ink"], font(24))
                    y += 120
            elif spec.template == "bullets_asset":
                draw.rounded_rectangle((82, 220, 665, 800), radius=26, fill=PALETTE["paper"], outline=PALETTE["line"], width=3)
                y = 265
                for bullet in spec.bullets or []:
                    draw.ellipse((122, y + 8, 136, y + 22), fill=PALETTE["teal"])
                    draw_wrapped_text(draw, (154, y, 620, y + 82), bullet, PALETTE["ink"], font(20))
                    y += 110
                pic = ASSETS / spec.asset if spec.asset and (ASSETS / spec.asset).exists() else REVIEWS / spec.asset
                panel = Image.open(pic).convert("RGB").resize((760, 560))
                img.paste(panel, (730, 228))
            elif spec.template == "full_asset":
                pic = ASSETS / spec.asset if spec.asset and (ASSETS / spec.asset).exists() else REVIEWS / spec.asset
                panel = Image.open(pic).convert("RGB").resize((1450, 600))
                img.paste(panel, (75, 220))
            elif spec.template == "table":
                preview_table(draw, (80, 220, 1520, 808), spec.table_headers or [], spec.table_rows or [])
            elif spec.template == "table_bullets":
                preview_table(draw, (82, 230, 655, 790), spec.table_rows[0], spec.table_rows[1:])
                draw.rounded_rectangle((710, 230, 1518, 790), radius=24, fill=PALETTE["paper"], outline=PALETTE["line"], width=3)
                y = 290
                for bullet in spec.bullets or []:
                    draw.ellipse((748, y + 8, 764, y + 24), fill=PALETTE["teal"])
                    draw_wrapped_text(draw, (790, y, 1460, y + 90), bullet, PALETTE["ink"], font(22))
                    y += 120
            elif spec.template == "image_grid":
                draw.rounded_rectangle((82, 228, 430, 800), radius=24, fill=PALETTE["paper"], outline=PALETTE["line"], width=3)
                y = 276
                for bullet in spec.bullets or []:
                    draw.ellipse((124, y + 8, 138, y + 22), fill=PALETTE["teal"])
                    draw_wrapped_text(draw, (158, y, 390, y + 88), bullet, PALETTE["ink"], font(18))
                    y += 118
                imgs = spec.assets or []
                layout = [(460, 230, 500, 275), (985, 230, 500, 275), (460, 520, 1025, 265)]
                for name, (x, y0, w, h) in zip(imgs, layout):
                    panel = Image.open(REVIEWS / name).convert("RGB").resize((w, h))
                    img.paste(panel, (x, y0))
            if spec.small_note:
                draw_wrapped_text(draw, (82, 830, 1520, 872), spec.small_note, PALETTE["muted"], font(13))
        img.save(PREVIEWS / f"slide_{idx:02d}.png")


def build_pdf(slide_count: int):
    c = canvas.Canvas(str(PDF_PATH), pagesize=(SLIDE_W * 72, SLIDE_H * 72))
    for idx in range(1, slide_count + 1):
        path = PREVIEWS / f"slide_{idx:02d}.png"
        c.drawImage(ImageReader(str(path)), 0, 0, width=SLIDE_W * 72, height=SLIDE_H * 72)
        c.showPage()
    c.save()


def write_notes(slides: list[SlideSpec]):
    lines = ["# DN 全量实验总结 deck notes", "", f"- 生成时间：{pd.Timestamp.now()}", f"- PPTX：`{PPTX_PATH}`", f"- PDF：`{PDF_PATH}`", ""]
    for idx, spec in enumerate(slides, start=1):
        lines.append(f"## Slide {idx:02d} · {spec.title}")
        lines.append(f"- section: {spec.section}")
        lines.append(f"- conclusion: {spec.conclusion}")
        if spec.bullets:
            lines.append("- bullets:")
            for bullet in spec.bullets:
                lines.append(f"  - {bullet}")
        if spec.small_note:
            lines.append(f"- note: {spec.small_note}")
        lines.append("")
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def main():
    ensure_dirs()
    build_assets()
    first_turn, next_turn, benchmark = load_data()
    slides = build_slide_specs(first_turn, next_turn, benchmark)
    build_ppt(slides)
    build_previews(slides)
    build_pdf(len(slides))
    write_notes(slides)
    print(json.dumps({
        "pptx": str(PPTX_PATH),
        "pdf": str(PDF_PATH),
        "preview_dir": str(PREVIEWS),
        "slide_count": len(slides),
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
