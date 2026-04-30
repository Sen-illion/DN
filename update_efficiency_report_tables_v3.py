from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pathlib import Path
from project_paths import path_in_project
from PIL import Image, ImageDraw, ImageFont
import zipfile, os

OUT = path_in_project('experiments', 'paper_method_view', '0_overview', 'efficiency_progress_deck_2026-04-26')
OUT.mkdir(parents=True, exist_ok=True)
pptx_path = OUT / 'DN_efficiency_progress_report_formal_v3_with_tables_2026-04-26.pptx'  # updated file
preview_dir = OUT / 'formal_v2_previews'
preview_dir.mkdir(exist_ok=True)

prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
W,H=prs.slide_width, prs.slide_height

NAVY=RGBColor(27,45,72); BLUE=RGBColor(39,95,153); TEAL=RGBColor(42,128,118); GREEN=RGBColor(61,139,92); ORANGE=RGBColor(213,128,56); RED=RGBColor(182,72,72); INK=RGBColor(35,39,47); MUTED=RGBColor(91,99,112); LIGHT=RGBColor(246,248,251); WHITE=RGBColor(255,255,255); LINE=RGBColor(214,221,230); PALEBLUE=RGBColor(238,244,250); PALEGREEN=RGBColor(236,246,243); PALEORANGE=RGBColor(255,248,238)
FONT='Microsoft YaHei'

def slide_bg(slide):
    bg=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,H); bg.fill.solid(); bg.fill.fore_color.rgb=LIGHT; bg.line.fill.background()
    top=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.11)); top.fill.solid(); top.fill.fore_color.rgb=NAVY; top.line.fill.background()

def tx(slide,text,x,y,w,h,size=12,color=INK,bold=False,align=PP_ALIGN.LEFT):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.02); tf.margin_right=Inches(0.02); tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01); tf.vertical_anchor=MSO_VERTICAL_ANCHOR.TOP
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return box

def para(slide,lines,x,y,w,h,size=9.5,color=MUTED,bullet=True):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.02); tf.margin_right=Inches(0.02); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=('• ' if bullet else '')+line; p.font.name=FONT; p.font.size=Pt(size); p.font.color.rgb=color
    return box

def title(slide,t,sub=None):
    tx(slide,'效率实验阶段性进展',0.55,0.34,3.0,0.24,10,BLUE,True)
    tx(slide,t,0.55,0.64,10.8,0.42,23.5,INK,True)
    if sub: tx(slide,sub,0.56,1.1,11.0,0.3,11.5,MUTED)
    line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0.55),Inches(1.52),Inches(12.7),Inches(1.52)); line.line.color.rgb=LINE; line.line.width=Pt(1)

def panel(slide,x,y,w,h,fill=WHITE,line=LINE):
    shp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.color.rgb=line; shp.line.width=Pt(0.8); return shp

def badge(slide,text,x,y,w,color,size=8.5):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.26)); shp.fill.solid(); shp.fill.fore_color.rgb=color; shp.line.fill.background()
    tf=shp.text_frame; tf.clear(); tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05); tf.margin_top=0; tf.margin_bottom=0; tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=WHITE

def metric(slide,label,value,note,x,y,w,color):
    tx(slide,value,x,y,w,0.35,21,color,True); tx(slide,label,x,y+0.4,w,0.2,9.5,INK,True); tx(slide,note,x,y+0.64,w,0.32,8.5,MUTED)

def table(slide,rows,x,y,w,h,widths=None,header=NAVY,fs=7.3,hs=7.5):
    shape=slide.shapes.add_table(len(rows),len(rows[0]),Inches(x),Inches(y),Inches(w),Inches(h)); tb=shape.table
    if widths:
        for i,cw in enumerate(widths): tb.columns[i].width=Inches(cw)
    for r,row in enumerate(rows):
        for c,val in enumerate(row):
            cell=tb.cell(r,c); cell.text=str(val); cell.margin_left=Inches(0.04); cell.margin_right=Inches(0.04); cell.margin_top=Inches(0.025); cell.margin_bottom=Inches(0.025); cell.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = header if r==0 else (WHITE if r%2 else RGBColor(241,245,249))
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER if c>0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name=FONT; run.font.size=Pt(hs if r==0 else fs); run.font.bold=(r==0); run.font.color.rgb=WHITE if r==0 else INK
    return shape

def hbar(slide,label,value,maxv,x,y,w,color,suffix='s',labw=1.0):
    tx(slide,label,x,y+0.02,labw,0.18,8.5,INK,False)
    barw=(w-labw-0.45)*value/maxv
    bg=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+labw),Inches(y+0.05),Inches(w-labw-0.45),Inches(0.13)); bg.fill.solid(); bg.fill.fore_color.rgb=RGBColor(229,234,240); bg.line.fill.background()
    b=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x+labw),Inches(y+0.05),Inches(max(0.02,barw)),Inches(0.13)); b.fill.solid(); b.fill.fore_color.rgb=color; b.line.fill.background()
    tx(slide,f'{value:.3f}{suffix}' if value<1 else f'{value:.2f}{suffix}',x+labw+barw+0.05,y-0.005,0.55,0.18,7.5,MUTED)

# Data
baseline_rows=[['System','Role','n','First mean','p95','Next','Success'],['DN','Ours / full chain',20,7.662,17.716,'-',1.0],['LIGHT','External interactive baseline',8,0.410,0.477,0.431,1.0],['PWR','Story speed/reference',5,0.833,0.963,0.837,1.0],['GenAgents','Continuity supplement',8,9.760,18.697,6.285,0.875]]
readwait=[(50,0.0,0.375,0.019,0.021),(55,0.0,0.375,0.018,0.018),(60,0.333,0.833,0.017,0.013),(65,0.0,0.375,0.022,0.021),(70,0.0,0.5,0.019,0.018)]

# Slide 1
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
tx(s,'DN 项目效率实验阶段性进展',0.62,0.62,7.7,0.55,28,NAVY,True); tx(s,'Efficiency Experiment Progress Report · 2026-04-26',0.64,1.12,5.2,0.25,11.5,MUTED); line=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0.62),Inches(1.55),Inches(6.2),Inches(1.55)); line.line.color.rgb=BLUE; line.line.width=Pt(2)
tx(s,'本阶段聚焦“效率”证据：端到端耗时、可玩内容返回延迟、外部系统对照与关键机制消融。生成质量、人评与图文一致性将在后续质量评估阶段合并。',0.64,1.83,11.8,0.58,14.5,INK)
for i,(lab,val,note,col) in enumerate([('DN 默认 fullchain','20/20','完整链路成功；可作为效率主行',GREEN),('外部 baseline','4 类','LIGHT / PWR / GenAgents / WorldGeneration',BLUE),('统一协议','已建立','playable-latency schema + subset v1/v2',ORANGE),('写作包','已落盘','主表、caption、局限性说明',RED)]):
    x=0.7+i*3.18; panel(s,x,2.86,2.82,1.42); metric(s,lab,val,note,x+0.22,3.1,2.3,col)
panel(s,0.72,5.04,11.95,1.25,PALEBLUE,RGBColor(194,208,226)); tx(s,'当前结论',0.96,5.24,1.4,0.25,13,BLUE,True); tx(s,'效率实验已经从“DN 自身验证”推进到“DN 与外部系统在统一 playable-latency 口径下的对比”。本阶段只讨论效率，不提前得出生成质量优劣结论。',2.05,5.18,9.8,0.52,16,INK,True); tx(s,'定位：效率实验第一版证据包，后续与质量评估合并成最终综合实验结论。',2.05,5.82,8.6,0.22,10.5,MUTED)

# Slide 2
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'实验口径：统一 playable-latency protocol','将不同系统统一转换为“从触发到返回可继续游玩内容”的效率指标。')
for i,(h,b,col) in enumerate([('输入集合','efficiency_playable_subset_v1/v2\n固定主题与触发输入',TEAL),('系统运行','DN 默认链路 + 外部 baseline adapter\n记录 first playable 与 next turn',BLUE),('汇总指标','mean / p95 / success_rate\ncompleteness / continuity 作效率护栏',ORANGE),('论文表格','核心主表 + WorldGeneration 补充表\n明确 baseline 角色边界',RED)]):
    x=0.75+i*3.12; panel(s,x,2.02,2.62,1.42); badge(s,str(i+1),x+0.12,2.18,0.36,col); tx(s,h,x+0.62,2.12,1.75,0.25,13.5,INK,True); tx(s,b,x+0.18,2.57,2.25,0.55,9.8,MUTED)
    if i<3:
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x+2.68),Inches(2.72),Inches(x+3.0),Inches(2.72)); ln.line.color.rgb=MUTED; ln.line.width=Pt(1.4); ln.line.end_arrowhead=True
rows=[['指标','含义','当前用途'],['first_playable_time_s','首次返回可继续游玩的内容所需时间','主效率指标'],['p95_latency_s','高延迟尾部表现','观察稳定性/长尾风险'],['next_turn_time_s','后续一轮动作后的可玩返回时间','补充交互连续性'],['success_rate','运行是否成功返回可用内容','效率比较前提'],['completeness / continuity','输出完整性与连续性代理指标','效率护栏，不替代质量评价']]
table(s,rows,0.8,4.05,7.25,2.32,[1.75,3.1,2.4],fs=8.0,hs=8.5)
panel(s,8.5,4.05,3.85,2.32,RGBColor(255,252,244),RGBColor(229,205,166)); tx(s,'口径边界',8.75,4.28,1.6,0.25,14,ORANGE,True); para(s,['DN 行使用 fullchain generate_option latency 作为 first-playable proxy。','不同 baseline 的系统形态不完全同构，因此主表解释重点是效率位置和系统角色。','生成质量、语义贴合、人评和图像一致性不在本阶段下结论。'],8.75,4.72,3.25,1.18,9.2,MUTED)

# Slide 3 DN detail + chart
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'DN 默认配置效率画像：主链路已具备可量化结果','直接放入 DN 主表数据，并用阶段耗时图展示延迟结构。')
# stage bars left
stages=[('worldview\ndefault',33.657,68.872,BLUE),('fullchain\nworldview',12.458,30.515,TEAL),('generate\noption',7.662,17.716,GREEN),('main\ncharacter',55.912,77.015,ORANGE)]
tx(s,'阶段 mean / p95 耗时',0.72,1.82,2.2,0.25,12.5,INK,True)
maxv=80
for i,(lab,mean,p95,col) in enumerate(stages):
    y=2.25+i*0.72; tx(s,lab,0.72,y-0.04,1.0,0.35,8.3,INK); 
    # p95 background
    p95w=3.8*p95/maxv; meanw=3.8*mean/maxv
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.75),Inches(y+0.02),Inches(max(0.02,p95w)),Inches(0.18)); bg.fill.solid(); bg.fill.fore_color.rgb=RGBColor(214,221,230); bg.line.fill.background()
    b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.75),Inches(y+0.02),Inches(max(0.02,meanw)),Inches(0.18)); b.fill.solid(); b.fill.fore_color.rgb=col; b.line.fill.background()
    tx(s,f'mean {mean:.1f}s / p95 {p95:.1f}s',1.75+p95w+0.08,y-0.02,1.5,0.18,7.6,MUTED)
# raw table right
rows=[['DN 阶段','样本','mean(s)','median(s)','p95(s)','说明'],['worldview default','20','33.657','23.664','68.872','独立世界观生成'],['fullchain worldview','20','12.458','9.003','30.515','完整链路中的世界观阶段'],['generate option','20','7.662','0.022','17.716','当前 first-playable proxy'],['main character','20','55.912','56.546','77.015','主角图完成耗时']]
table(s,rows,6.0,1.95,6.5,2.5,[1.55,0.55,0.76,0.8,0.73,2.1],fs=7.1,hs=7.4)
# guardrails direct table + text
rows2=[['有效性代理','结果'],['worldview_success_rate','1.0'],['first_scene_success_rate','1.0'],['image_return_rate','1.0'],['option_count_ge_2_rate','1.0'],['fallback_trigger_rate','0.0'],['scene_prompt_pollution_rate','0.0']]
table(s,rows2,0.78,5.28,4.18,1.35,[2.95,0.95],fs=7.0,hs=7.6,header=TEAL)
panel(s,5.28,5.28,7.12,1.35,PALEBLUE,RGBColor(194,208,226)); tx(s,'解读',5.55,5.47,0.75,0.22,12,BLUE,True); tx(s,'DN 默认配置已经具备可量化的完整链路效率画像；first-playable proxy 可用于外部系统效率对比，但质量判断仍需后续人评 / judge / 图文一致性评估。',6.2,5.42,5.75,0.45,10.5,INK,True)

# Slide 4 baseline status table
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'外部 baseline 已补齐：每个系统有明确角色与证据边界','这里将已完成的 baseline 状态表直接放入 PPT，避免只用文字概述。')
rows=[['Baseline','当前角色','样本','已完成证据','主要限制'],['LIGHT','权威互动对话 / game-world baseline','8','官方 checkpoint；8-item playable latency batch','英文 adapter；语义贴合 DN 中文主题有限'],['Plan-Write-Revise','on-demand story generation 速度参考','5','官方模型包；5-item batch；summary 已生成','轻量文本系统，非完整多模态游戏链路'],['GenAgents','状态连续性 / 多轮 agent 补充 baseline','8','8-item live run 转换为 playable-latency summary','不输出 DN-style 图像、世界观 JSON 或分支系统'],['WorldGeneration','世界构造补充 baseline','8','fallback graph-to-playable 路径；8-item summary','非完整原论文 pipeline，建议放补充表'],['AIDungeon / StoryDiffusion','当前周期排除或延期','-','status / decision note 已记录','legacy runtime 或 CUDA 硬件阻塞']]
table(s,rows,0.62,1.85,12.1,3.62,[1.75,2.2,0.55,3.1,4.5],fs=7.0,hs=7.5)
for i,(h,t,c) in enumerate([('核心主表','DN + LIGHT + PWR + GenAgents',BLUE),('补充说明','WorldGeneration fallback row',ORANGE),('当前不进入','AIDungeon / StoryDiffusion',RED)]):
    x=0.85+i*4.1; panel(s,x,5.92,3.45,0.72); badge(s,h,x+0.15,6.12,1.3,c,8); tx(s,t,x+1.58,6.1,1.7,0.28,8.8,MUTED)
tx(s,'整理原则：只有具备 source_links + protocol + raw_runs + summaries + status 的 baseline，才进入主表或补充表讨论。',0.75,6.9,11.7,0.22,9.0,MUTED)

# Slide 5 main table + graphics
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'主表数据与图形化展示：DN 与外部系统 first-playable 对比','左侧为主表关键数据，右侧用条形图展示 first mean 与 p95 的相对位置。')
rows=[['System','Role','n','First mean','p95','Next','Success'],['DN','Ours / full chain','20','7.662','17.716','-','1.0'],['LIGHT','External interactive','8','0.410','0.477','0.431','1.0'],['PWR','Story speed/ref.','5','0.833','0.963','0.837','1.0'],['GenAgents','Continuity supplement','8','9.760','18.697','6.285','0.875']]
table(s,rows,0.72,1.95,6.1,2.25,[1.0,1.55,0.38,0.78,0.7,0.62,0.65],fs=7.2,hs=7.4,header=BLUE)
tx(s,'First playable mean (s)',7.15,1.92,2.5,0.25,11.2,INK,True)
for i,(lab,val,col) in enumerate([('LIGHT',0.410,BLUE),('PWR',0.833,ORANGE),('DN',7.662,GREEN),('GenAgents',9.760,RED)]): hbar(s,lab,val,10,7.15,2.32+i*0.48,4.4,col)
tx(s,'p95 latency (s)',7.15,4.38,2.0,0.25,11.2,INK,True)
for i,(lab,val,col) in enumerate([('LIGHT',0.477,BLUE),('PWR',0.963,ORANGE),('DN',17.716,GREEN),('GenAgents',18.697,RED)]): hbar(s,lab,val,20,7.15,4.72+i*0.34,4.4,col)
# notes bottom
for x,h,b,c in [(0.75,'效率层面','外部轻量/对话系统显著更快；DN 完成的是更复杂的完整链路。',TEAL),(4.75,'边界','不能直接由速度表推出生成质量优劣，质量评估后续补充。',ORANGE),(8.75,'写法','表述为“统一口径下的效率位置”，避免写成全面胜负。',BLUE)]:
    panel(s,x,6.18,3.55,0.52,WHITE,LINE); tx(s,h,x+0.18,6.31,1.0,0.18,9.8,c,True); tx(s,b,x+0.95,6.29,2.35,0.22,7.2,MUTED)

# Slide 6 readwait threshold table + chart
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'预生成 / readwait 阈值数据：60s 是当前最稳定收益点','直接放入 50/55/60/65/70s 阈值表，并用图形展示 real-scene 命中率变化。')
rows=[['read wait','off real rate','on real rate','off median','on median','样本'],['50s','0.000','0.375','0.019','0.021','8v8'],['55s','0.000','0.375','0.018','0.018','8v8'],['60s','0.333','0.833','0.017','0.013','12v12'],['65s','0.000','0.375','0.022','0.021','8v8'],['70s','0.000','0.500','0.019','0.018','7v8 / 8v8']]
table(s,rows,0.7,1.95,5.55,2.75,[0.8,0.9,0.9,0.85,0.85,1.25],fs=7.3,hs=7.5,header=TEAL)
# line chart for real rate
chart_data=CategoryChartData(); chart_data.categories=['50s','55s','60s','65s','70s']; chart_data.add_series('pregen off real-scene rate',(0,0,0.333,0,0)); chart_data.add_series('pregen on real-scene rate',(0.375,0.375,0.833,0.375,0.5))
chart=s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,Inches(6.75),Inches(1.95),Inches(5.3),Inches(2.75),chart_data).chart
chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout=False
chart.value_axis.maximum_scale=1.0; chart.value_axis.minimum_scale=0; chart.value_axis.tick_labels.font.size=Pt(7); chart.category_axis.tick_labels.font.size=Pt(8)
# median bars small using direct labels
panel(s,0.7,5.18,5.55,1.18,PALEGREEN,RGBColor(194,221,214)); tx(s,'关键主结果（60s）',0.95,5.38,1.8,0.22,12.5,TEAL,True); tx(s,'真实剧情命中率：33.3% → 83.3%；第二次点击中位延迟：0.017s → 0.013s。说明预生成的收益主要体现在“下一次交互”而不是当前请求无条件更快。',0.95,5.75,4.85,0.35,9.8,INK,True)
panel(s,6.75,5.18,5.3,1.18,PALEORANGE,RGBColor(229,205,166)); tx(s,'解释边界',7.0,5.38,1.5,0.22,12.5,ORANGE,True); para(s,['real_scene 当前为启发式近似。','除 60s 外多为 8v8，适合阈值趋势说明。','70s off 已剔除息屏异常样本。'],7.0,5.72,4.2,0.45,8.5,MUTED)

# Slide 7 mechanism summary table + visuals
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'内部机制效率实验：已有结果解释 DN 延迟结构','该页把机制消融的关键数字放进表格，并图形化展示其含义。')
rows=[['机制','对比组 / 口径','关键数据','效率含义'],['Readwait / Pregeneration','60s next-click','real-scene 33.3%→83.3%；median 0.017s→0.013s','预生成收益体现为下一次点击体验'],['Council','default vs no-council fullchain','default 20/20；no-council 19/20；worldview mean 12.458s vs 16.547s','影响稳定性与延迟结构'],['Pregen clean12','pregen off vs on','worldview mean 73.377s→16.566s','对主链路延迟控制有贡献']]
table(s,rows,0.72,1.95,7.25,2.55,[1.25,1.55,2.35,2.1],fs=7.3,hs=7.6,header=ORANGE)
# right mini bars for pregen clean
panel(s,8.35,1.95,3.75,2.55,WHITE,LINE); tx(s,'Pregen clean12 worldview mean',8.6,2.18,2.5,0.22,11,INK,True); hbar(s,'off',73.377,80,8.6,2.75,2.9,RED,labw=0.45); hbar(s,'on',16.566,80,8.6,3.25,2.9,GREEN,labw=0.45); tx(s,'从 73.377s 降至 16.566s；但机制解释需结合 readwait next-click 口径。',8.6,3.88,3.05,0.42,8.8,MUTED)
# bottom narrative
panel(s,0.75,5.25,11.5,1.15,PALEBLUE,RGBColor(194,208,226)); tx(s,'本阶段可写入论文/汇报的效率结论',1.0,5.45,2.6,0.22,12.5,BLUE,True); tx(s,'DN 的效率证据目前由三部分构成：默认链路效率画像、外部系统 playable-latency 对比、关键机制消融。完整的“效率-质量权衡”需要等待后续质量评估结果合并，当前不建议单独扩展完整参数矩阵。',1.0,5.85,10.55,0.35,10.6,INK,True)

# Slide 8 deliverables and next steps
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s); title(s,'当前可交付材料与后续工作','收束本阶段产物：哪些材料已经形成 source of truth，哪些工作留给下一阶段。')
panel(s,0.75,1.9,5.55,3.9); tx(s,'已形成的材料',1.02,2.15,2.0,0.28,15,TEAL,True); para(s,['主表：main_playable_latency_scaffold_2026-04-26.csv','补充表：main_playable_latency_with_light；supplementary_worldgeneration','协议与子集：schema、playable_protocol、subset v1/v2','写作包：实验章节草稿、caption、limitations、bundle index','baseline 均有 status、raw_runs、summaries 或 decision note'],1.02,2.62,4.95,2.6,9.0,MUTED)
panel(s,6.75,1.9,5.55,3.9); tx(s,'下一阶段建议',7.02,2.15,2.0,0.28,15,BLUE,True); para(s,['冻结效率主表口径：DN / LIGHT / PWR / GenAgents 核心主表，WorldGeneration 补充。','修正文稿表达：避免“全面优于 baseline”等不严谨说法。','等待质量组补充文本质量、人评、图文一致性、视觉质量。','完整 DN 参数矩阵暂不作为当前效率主线，因为只看速度会误导。'],7.02,2.62,4.95,2.4,9.0,MUTED)
panel(s,0.75,6.18,11.55,0.62,RGBColor(244,247,251),RGBColor(203,213,225)); tx(s,'阶段性结论：效率实验第一版证据包已经完成，现阶段重点是冻结表格口径与规范解释；质量评估完成后，再形成最终综合实验结论。',1.0,6.36,10.95,0.25,11.0,INK,True)

prs.save(pptx_path)
with zipfile.ZipFile(pptx_path,'r') as z:
    slides=[n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
    assert len(slides)==8

# Previews (approximate but includes updated slide gist)
try:
    F_TITLE=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc',38); F_H1=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc',30); F_H2=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc',22); F_BODY=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc',17); F_SMALL=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc',13)
except: F_TITLE=F_H1=F_H2=F_BODY=F_SMALL=ImageFont.load_default()
PW,PH=1600,900; C_LIGHT=(246,248,251); C_NAVY=(27,45,72); C_BLUE=(39,95,153); C_TEAL=(42,128,118); C_GREEN=(61,139,92); C_ORANGE=(213,128,56); C_RED=(182,72,72); C_INK=(35,39,47); C_MUTED=(91,99,112); C_LINE=(214,221,230); C_WHITE=(255,255,255)

def wrap(d,text,x,y,font,fill,maxw,gap=4):
    for para in text.split('\n'):
        line=''
        for ch in para:
            if d.textbbox((0,0),line+ch,font=font)[2]<=maxw or not line: line+=ch
            else: d.text((x,y),line,font=font,fill=fill); y+=font.size+gap; line=ch
        if line: d.text((x,y),line,font=font,fill=fill); y+=font.size+gap
    return y

def base(t,sub=''):
    img=Image.new('RGB',(PW,PH),C_LIGHT); d=ImageDraw.Draw(img); d.rectangle([0,0,PW,13],fill=C_NAVY); d.text((66,38),'效率实验阶段性进展',font=F_SMALL,fill=C_BLUE); d.text((66,76),t,font=F_H1,fill=C_INK); 
    if sub: d.text((66,130),sub,font=F_SMALL,fill=C_MUTED)
    d.line([66,180,1530,180],fill=C_LINE,width=2); return img,d

def dt(d,rows,x,y,w,h,widths=None,header=C_NAVY):
    n=len(rows); rh=h/n; widths=widths or [w/len(rows[0])]*len(rows[0])
    for r,row in enumerate(rows):
        cx=x
        for c,val in enumerate(row):
            cw=widths[c]; fill=header if r==0 else (C_WHITE if r%2 else (241,245,249)); d.rectangle([cx,y+r*rh,cx+cw,y+(r+1)*rh],fill=fill,outline=C_LINE); wrap(d,str(val),cx+6,y+r*rh+6,F_SMALL,(255,255,255) if r==0 else C_INK,cw-12,1); cx+=cw

def prev1():
    img=Image.new('RGB',(PW,PH),C_LIGHT); d=ImageDraw.Draw(img); d.rectangle([0,0,PW,13],fill=C_NAVY); d.text((75,72),'DN 项目效率实验阶段性进展',font=F_TITLE,fill=C_NAVY); d.text((76,130),'Efficiency Experiment Progress Report · 2026-04-26',font=F_BODY,fill=C_MUTED); d.line([76,180,750,180],fill=C_BLUE,width=4); wrap(d,'本阶段聚焦“效率”证据：端到端耗时、可玩内容返回延迟、外部系统对照与关键机制消融。生成质量、人评与图文一致性将在后续质量评估阶段合并。',76,215,F_BODY,C_INK,1380)
    for i,(lab,val,note,c) in enumerate([('DN 默认 fullchain','20/20','完整链路成功；可作为效率主行',C_GREEN),('外部 baseline','4 类','LIGHT / PWR / GenAgents / WorldGeneration',C_BLUE),('统一协议','已建立','playable-latency schema + subset v1/v2',C_ORANGE),('写作包','已落盘','主表、caption、局限性说明',C_RED)]):
        x=85+i*385; d.rectangle([x,340,x+340,510],fill=C_WHITE,outline=C_LINE,width=2); d.text((x+25,365),val,font=F_TITLE,fill=c); d.text((x+25,425),lab,font=F_BODY,fill=C_INK); wrap(d,note,x+25,455,F_SMALL,C_MUTED,290)
    d.rectangle([85,610,1510,760],fill=(238,244,250),outline=(194,208,226)); wrap(d,'当前结论：效率实验已经从“DN 自身验证”推进到“DN 与外部系统在统一 playable-latency 口径下的对比”。本阶段只讨论效率，不提前得出生成质量优劣结论。',120,635,F_H2,C_INK,1300); return img

def prev2():
    img,d=base('实验口径：统一 playable-latency protocol','将不同系统统一转换为“从触发到返回可继续游玩内容”的效率指标。')
    for i,txt in enumerate(['输入集合\nefficiency subset v1/v2','系统运行\nDN + baseline adapters','汇总指标\nmean / p95 / success','论文表格\n核心主表 + 补充表']):
        x=90+i*370; d.rectangle([x,250,x+310,395],fill=C_WHITE,outline=C_LINE,width=2); d.ellipse([x+15,270,x+55,310],fill=[C_TEAL,C_BLUE,C_ORANGE,C_RED][i]); d.text((x+28,278),str(i+1),font=F_SMALL,fill='white'); wrap(d,txt,x+75,268,F_BODY,C_INK,210); 
        if i<3: d.line([x+318,322,x+360,322],fill=C_MUTED,width=3)
    dt(d,[['指标','含义','当前用途'],['first_playable_time_s','首次返回可继续游玩的内容所需时间','主效率指标'],['p95_latency_s','高延迟尾部表现','长尾风险'],['next_turn_time_s','后续一轮动作后的返回时间','交互连续性'],['success_rate','是否成功返回可用内容','比较前提'],['completeness / continuity','输出完整性与连续性代理','效率护栏']],95,490,870,280,[210,360,300])
    d.rectangle([1015,490,1480,770],fill=(255,252,244),outline=(229,205,166)); wrap(d,'口径边界\n• DN 行使用 generate_option latency 作为 proxy。\n• baseline 不是完全同构系统。\n• 质量评估后续补充。',1045,520,F_SMALL,C_MUTED,380); return img

def prev3():
    img,d=base('DN 默认配置效率画像：主链路已具备可量化结果','直接放入 DN 主表数据，并用阶段耗时图展示延迟结构。')
    maxv=80; stages=[('worldview default',33.657,68.872,C_BLUE),('fullchain worldview',12.458,30.515,C_TEAL),('generate option',7.662,17.716,C_GREEN),('main character',55.912,77.015,C_ORANGE)]
    d.text((85,225),'阶段 mean / p95 耗时',font=F_BODY,fill=C_INK)
    for i,(lab,mean,p95,c) in enumerate(stages):
        y=280+i*85; d.text((85,y),lab,font=F_SMALL,fill=C_INK); d.rectangle([260,y+5,260+int(p95/maxv*450),y+24],fill=C_LINE); d.rectangle([260,y+5,260+int(mean/maxv*450),y+24],fill=c); d.text((730,y),f'mean {mean:.1f}s / p95 {p95:.1f}s',font=F_SMALL,fill=C_MUTED)
    dt(d,[['DN 阶段','样本','mean','median','p95','说明'],['worldview default','20','33.657','23.664','68.872','独立世界观生成'],['fullchain worldview','20','12.458','9.003','30.515','完整链路世界观'],['generate option','20','7.662','0.022','17.716','first-playable proxy'],['main character','20','55.912','56.546','77.015','主角图完成耗时']],720,235,780,300,[180,65,85,85,85,280])
    dt(d,[['有效性代理','结果'],['worldview_success_rate','1.0'],['first_scene_success_rate','1.0'],['image_return_rate','1.0'],['option_count_ge_2_rate','1.0'],['fallback_trigger_rate','0.0'],['scene_prompt_pollution_rate','0.0']],90,635,500,160,[360,140],header=C_TEAL)
    d.rectangle([650,635,1500,790],fill=(238,244,250),outline=(194,208,226)); wrap(d,'解读：DN 默认配置已经具备可量化的完整链路效率画像；first-playable proxy 可用于外部系统效率对比，但质量判断仍需后续人评 / judge / 图文一致性评估。',680,675,F_BODY,C_INK,770); return img

def prev4():
    img,d=base('外部 baseline 已补齐：每个系统有明确角色与证据边界','这里将已完成的 baseline 状态表直接放入 PPT。')
    dt(d,[['Baseline','当前角色','样本','已完成证据','主要限制'],['LIGHT','权威互动对话 / game-world','8','官方 checkpoint；8-item batch','英文 adapter；语义贴合 DN 中文主题有限'],['Plan-Write-Revise','on-demand story 速度参考','5','官方模型包；5-item summary','轻量文本系统，非完整多模态链路'],['GenAgents','状态连续性 / 多轮 agent','8','8-item live run 转换 summary','不输出 DN-style 图像、世界观 JSON'],['WorldGeneration','世界构造补充 baseline','8','fallback graph-to-playable summary','非完整原论文 pipeline'],['AIDungeon / StoryDiffusion','当前周期排除或延期','-','status / decision note','legacy runtime 或 CUDA 阻塞']],75,220,1450,430,[185,270,70,395,530])
    for i,(h,t,c) in enumerate([('核心主表','DN + LIGHT + PWR + GenAgents',C_BLUE),('补充说明','WorldGeneration fallback row',C_ORANGE),('当前不进入','AIDungeon / StoryDiffusion',C_RED)]):
        x=100+i*490; d.rectangle([x,705,x+390,780],fill=C_WHITE,outline=C_LINE); d.rounded_rectangle([x+15,725,x+135,755],radius=12,fill=c); d.text((x+28,731),h,font=F_SMALL,fill='white'); d.text((x+155,728),t,font=F_SMALL,fill=C_MUTED)
    return img

def prev5():
    img,d=base('主表数据与图形化展示：DN 与外部系统 first-playable 对比','左侧为主表关键数据，右侧用条形图展示 first mean 与 p95 的相对位置。')
    dt(d,[['System','Role','n','First','p95','Next','Success'],['DN','Ours / full chain','20','7.662','17.716','-','1.0'],['LIGHT','External interactive','8','0.410','0.477','0.431','1.0'],['PWR','Story speed/ref.','5','0.833','0.963','0.837','1.0'],['GenAgents','Continuity supplement','8','9.760','18.697','6.285','0.875']],85,235,730,270,[110,190,50,90,80,70,85],header=C_BLUE)
    d.text((860,230),'First playable mean (s)',font=F_BODY,fill=C_INK)
    for i,(lab,val,c) in enumerate([('LIGHT',0.410,C_BLUE),('PWR',0.833,C_ORANGE),('DN',7.662,C_GREEN),('GenAgents',9.760,C_RED)]):
        y=285+i*55; d.text((860,y),lab,font=F_SMALL,fill=C_INK); d.rectangle([990,y+8,1430,y+24],fill=C_LINE); d.rectangle([990,y+8,990+int(val/10*440),y+24],fill=c); d.text((1000+int(val/10*440),y),f'{val:.3f}s',font=F_SMALL,fill=C_MUTED)
    d.text((860,515),'p95 latency (s)',font=F_BODY,fill=C_INK)
    for i,(lab,val,c) in enumerate([('LIGHT',0.477,C_BLUE),('PWR',0.963,C_ORANGE),('DN',17.716,C_GREEN),('GenAgents',18.697,C_RED)]):
        y=560+i*38; d.text((860,y),lab,font=F_SMALL,fill=C_INK); d.rectangle([990,y+8,1430,y+24],fill=C_LINE); d.rectangle([990,y+8,990+int(val/20*440),y+24],fill=c); d.text((1000+int(val/20*440),y),f'{val:.3f}s',font=F_SMALL,fill=C_MUTED)
    for x,h,b,c in [(90,'效率层面','外部轻量/对话系统显著更快；DN 完成的是更复杂的完整链路。',C_TEAL),(570,'边界','不能由速度表推出质量优劣，质量评估后续补充。',C_ORANGE),(1050,'写法','表述为“统一口径下的效率位置”。',C_BLUE)]:
        d.rectangle([x,735,x+400,800],fill=C_WHITE,outline=C_LINE); d.text((x+15,748),h,font=F_SMALL,fill=c); wrap(d,b,x+100,748,F_SMALL,C_MUTED,275)
    return img

def prev6():
    img,d=base('预生成 / readwait 阈值数据：60s 是当前最稳定收益点','直接放入 50/55/60/65/70s 阈值表，并用图形展示 real-scene 命中率变化。')
    dt(d,[['read wait','off real','on real','off med','on med','样本'],['50s','0.000','0.375','0.019','0.021','8v8'],['55s','0.000','0.375','0.018','0.018','8v8'],['60s','0.333','0.833','0.017','0.013','12v12'],['65s','0.000','0.375','0.022','0.021','8v8'],['70s','0.000','0.500','0.019','0.018','7v8 / 8v8']],85,235,665,330,[95,105,105,100,100,160],header=C_TEAL)
    # simple line chart
    x0,y0=860,585; chartw, charth=550,300; d.rectangle([820,230,1490,585],fill=C_WHITE,outline=C_LINE); d.text((850,255),'real-scene rate',font=F_BODY,fill=C_INK); d.line([x0,y0,x0+chartw,y0],fill=C_LINE,width=2); d.line([x0,y0-charth,x0,y0],fill=C_LINE,width=2)
    vals_off=[0,0,0.333,0,0]; vals_on=[0.375,0.375,0.833,0.375,0.5]; xs=[x0+i*chartw/4 for i in range(5)]
    for vals,c in [(vals_off,C_RED),(vals_on,C_GREEN)]:
        pts=[(xs[i],y0-vals[i]*charth) for i in range(5)]; d.line(pts,fill=c,width=4)
        for x,y in pts: d.ellipse([x-5,y-5,x+5,y+5],fill=c)
    for i,l in enumerate(['50','55','60','65','70']): d.text((xs[i]-12,y0+12),l+'s',font=F_SMALL,fill=C_MUTED)
    d.text((1230,255),'on',font=F_SMALL,fill=C_GREEN); d.text((1280,255),'off',font=F_SMALL,fill=C_RED)
    d.rectangle([85,630,750,775],fill=(236,246,243),outline=(194,221,214)); wrap(d,'关键主结果（60s）：真实剧情命中率 33.3% → 83.3%；第二次点击中位延迟 0.017s → 0.013s。说明收益主要体现在“下一次交互”。',115,655,F_BODY,C_INK,590)
    d.rectangle([820,630,1490,775],fill=(255,248,238),outline=(229,205,166)); wrap(d,'解释边界：real_scene 当前为启发式近似；除 60s 外多为 8v8，适合阈值趋势说明；70s off 已剔除息屏异常样本。',850,655,F_BODY,C_MUTED,590)
    return img

def prev7():
    img,d=base('内部机制效率实验：已有结果解释 DN 延迟结构','机制消融关键数字表 + 图形化解释。')
    dt(d,[['机制','对比组 / 口径','关键数据','效率含义'],['Readwait / Pregeneration','60s next-click','real-scene 33.3%→83.3%；median 0.017s→0.013s','预生成收益体现为下一次点击体验'],['Council','default vs no-council fullchain','default 20/20；no-council 19/20；worldview mean 12.458s vs 16.547s','影响稳定性与延迟结构'],['Pregen clean12','pregen off vs on','worldview mean 73.377s→16.566s','对主链路延迟控制有贡献']],85,235,870,300,[150,190,280,250],header=C_ORANGE)
    d.rectangle([1000,235,1450,530],fill=C_WHITE,outline=C_LINE); d.text((1030,260),'Pregen clean12 worldview mean',font=F_BODY,fill=C_INK)
    for i,(lab,val,c) in enumerate([('off',73.377,C_RED),('on',16.566,C_GREEN)]):
        y=330+i*70; d.text((1030,y),lab,font=F_SMALL,fill=C_INK); d.rectangle([1090,y+8,1390,y+24],fill=C_LINE); d.rectangle([1090,y+8,1090+int(val/80*300),y+24],fill=c); d.text((1100+int(val/80*300),y),f'{val:.1f}s',font=F_SMALL,fill=C_MUTED)
    d.rectangle([85,650,1490,780],fill=(238,244,250),outline=(194,208,226)); wrap(d,'本阶段可写入论文/汇报的效率结论：DN 的效率证据目前由三部分构成——默认链路效率画像、外部系统 playable-latency 对比、关键机制消融。完整的“效率-质量权衡”需要等待后续质量评估结果合并。',115,680,F_BODY,C_INK,1320)
    return img

def prev8():
    img,d=base('当前可交付材料与后续工作','收束本阶段产物：哪些材料已经形成 source of truth，哪些工作留给下一阶段。')
    d.rectangle([90,230,750,690],fill=C_WHITE,outline=C_LINE); d.text((120,260),'已形成的材料',font=F_H2,fill=C_TEAL); wrap(d,'• 主表：main_playable_latency_scaffold_2026-04-26.csv\n• 补充表：main_playable_latency_with_light；supplementary_worldgeneration\n• 协议与子集：schema、playable_protocol、subset v1/v2\n• 写作包：实验章节草稿、caption、limitations、bundle index\n• baseline 均有 status、raw_runs、summaries 或 decision note。',120,315,F_SMALL,C_MUTED,570)
    d.rectangle([850,230,1510,690],fill=C_WHITE,outline=C_LINE); d.text((880,260),'下一阶段建议',font=F_H2,fill=C_BLUE); wrap(d,'• 冻结效率主表口径：DN / LIGHT / PWR / GenAgents 核心主表，WorldGeneration 补充。\n• 修正文稿表达：避免“全面优于 baseline”等不严谨说法。\n• 等待质量组补充文本质量、人评、图文一致性、视觉质量。\n• 完整 DN 参数矩阵暂不作为当前效率主线。',880,315,F_SMALL,C_MUTED,570)
    d.rectangle([90,740,1510,805],fill=(244,247,251),outline=(203,213,225)); wrap(d,'阶段性结论：效率实验第一版证据包已经完成，现阶段重点是冻结表格口径与规范解释；质量评估完成后，再形成最终综合实验结论。',120,757,F_BODY,C_INK,1320)
    return img

for old in preview_dir.glob('slide_*.png'): old.unlink()
imgs=[prev1(),prev2(),prev3(),prev4(),prev5(),prev6(),prev7(),prev8()]
paths=[]
for i,img in enumerate(imgs,1): p=preview_dir/f'slide_{i}.png'; img.save(p); paths.append(p)
thumbs=[Image.open(p).resize((360,203)) for p in paths]; rows=(len(thumbs)+1)//2; mont=Image.new('RGB',(760,rows*223),(235,238,242))
for i,img in enumerate(thumbs): mont.paste(img,(10+(i%2)*380,10+(i//2)*223))
montage=preview_dir/'montage.png'; mont.save(montage)
print('PPTX',pptx_path); print('slides',len(slides)); print('montage',montage)
