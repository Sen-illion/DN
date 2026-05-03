from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path
from project_paths import path_in_project
from PIL import Image, ImageDraw, ImageFont
import zipfile, os, math

OUT = path_in_project('experiments', 'paper_method_view', '0_overview', 'teacher_report_efficiency_2026-04-26')
OUT.mkdir(parents=True, exist_ok=True)
pptx_path = OUT / 'DN_efficiency_progress_report_2026-04-26.pptx'
preview_dir = OUT / 'previews'
preview_dir.mkdir(exist_ok=True)

# Colors
INK = RGBColor(35, 39, 47)
MUTED = RGBColor(95, 101, 113)
BG = RGBColor(248, 244, 235)
CREAM = RGBColor(255, 251, 242)
GREEN = RGBColor(56, 132, 103)
BLUE = RGBColor(54, 103, 161)
ORANGE = RGBColor(218, 124, 54)
RED = RGBColor(178, 73, 70)
YELLOW = RGBColor(232, 185, 83)
LINE = RGBColor(213, 205, 188)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helpers
W, H = prs.slide_width, prs.slide_height

def add_bg(slide, color=BG):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def tx(slide, text, x, y, w, h, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box

def pill(slide, text, x, y, w, h, fill, color=RGBColor(255,255,255), size=15):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.13); tf.margin_right = Inches(0.13); tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name='Microsoft YaHei'; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=color
    return shp

def rule(slide, x, y, w, color=INK, weight=2):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln

def card(slide, x,y,w,h, fill=CREAM, line=LINE, radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    shp.line.color.rgb=line; shp.line.width=Pt(1)
    return shp

def add_bullet(slide, text, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=0; tf.margin_top=0; tf.margin_bottom=0; tf.margin_right=0
    lines = text.split('\n')
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.level=0; p.text = line
        p.font.name='Microsoft YaHei'; p.font.size=Pt(size); p.font.color.rgb=color
    return box

# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
# big left accent
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), H)
accent.fill.solid(); accent.fill.fore_color.rgb = GREEN; accent.line.fill.background()
tx(s, 'DN 效率实验阶段性汇报', 0.75, 0.72, 7.2, 0.55, 26, GREEN, True)
tx(s, '已从“DN 自身跑通”\n推进到“外部 baseline 可玩延迟对比”', 0.72, 1.55, 8.9, 1.65, 44, INK, True)
tx(s, '当前重点：只汇报效率，不提前下生成质量结论；质量评估后续与组内成员共同完成。', 0.78, 3.35, 8.3, 0.55, 20, MUTED)
# metric typographic cluster
pill(s, 'DN fullchain 20/20 成功', 0.8, 4.55, 2.55, 0.55, GREEN, size=15)
pill(s, 'LIGHT 8-item 已跑通', 3.55, 4.55, 2.45, 0.55, BLUE, size=15)
pill(s, 'PWR / GenAgents 已入表', 6.2, 4.55, 2.75, 0.55, ORANGE, size=15)
# visual bars
for i,(label,val,col) in enumerate([('DN',7.662,GREEN),('PWR',0.833,ORANGE),('LIGHT',0.410,BLUE),('GenAgents',9.760,RED)]):
    x=9.6+i*0.72
    h=max(0.18, val/10*2.15)
    shp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.9-h), Inches(0.42), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=col; shp.line.fill.background()
    tx(s,label,x-0.08,6.02,0.7,0.28,10,MUTED,False,PP_ALIGN.CENTER)
    tx(s,f'{val:.1f}s' if val>=1 else f'{val:.2f}s',x-0.13,5.62-h,0.8,0.25,10,INK,True,PP_ALIGN.CENTER)
tx(s,'first playable mean',9.45,3.28,2.7,0.35,14,MUTED,True)
rule(s,9.45,3.68,2.3,GREEN,3)
tx(s,'汇报口径：统一 playable-latency protocol 下的效率证据包',0.78,6.9,8.5,0.28,11,MUTED)

# Slide 2: map
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tx(s,'现在已经做出的效率实验包',0.65,0.45,7.8,0.55,32,INK,True)
tx(s,'从默认链路、外部 baseline、机制消融到写作索引，已经形成可汇报材料。',0.67,1.0,9.0,0.38,17,MUTED)
# process nodes
nodes=[('1', 'DN 默认链路', '20/20 成功\n效率画像完成', GREEN),('2','统一协议','playable latency\nsubset v1/v2', BLUE),('3','外部 baseline','LIGHT / PWR /\nGenAgents / WG', ORANGE),('4','写作包','Table 1 + caption\n局限性说明', RED)]
for i,(num,title,body,col) in enumerate(nodes):
    x=0.8+i*3.1
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.05), Inches(0.68), Inches(0.68))
    circ.fill.solid(); circ.fill.fore_color.rgb=col; circ.line.fill.background()
    tf=circ.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=num; r.font.name='Microsoft YaHei'; r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)
    tx(s,title,x-0.15,2.95,2.35,0.32,19,INK,True)
    tx(s,body,x-0.15,3.4,2.25,0.62,15,MUTED)
    if i<3:
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+0.82), Inches(2.39), Inches(x+2.68), Inches(2.39))
        ln.line.color.rgb=LINE; ln.line.width=Pt(2.5)
        ln.line.end_arrowhead = True
# status lanes
card(s,0.75,4.7,3.4,1.3); tx(s,'已完成',0.95,4.88,1.4,0.28,18,GREEN,True); tx(s,'DN 主链路、readwait/预生成、Council 部分消融、四类 baseline 延迟结果',0.95,5.28,2.85,0.58,13,MUTED)
card(s,4.75,4.7,3.4,1.3); tx(s,'可汇报',4.95,4.88,1.4,0.28,18,BLUE,True); tx(s,'主表候选、补充表、caption、局限性与汇报段落已经落盘',4.95,5.28,2.85,0.58,13,MUTED)
card(s,8.75,4.7,3.4,1.3); tx(s,'暂不做',8.95,4.88,1.4,0.28,18,ORANGE,True); tx(s,'完整参数矩阵/质量优劣结论：等后续质量评估合并后再判断',8.95,5.28,2.85,0.58,13,MUTED)
tx(s,'关键 source of truth：main_experiment_writing_bundle_index_2026-04-26.md',0.78,6.9,9.0,0.26,11,MUTED)

# Slide 3 chart
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tx(s,'主实验表已补上外部系统效率对比',0.65,0.45,8.7,0.55,31,INK,True)
tx(s,'指标：first_playable_time_mean_s（从触发到返回可继续游玩内容）。数值越低越快，但系统能力范围不同。',0.67,0.98,10.8,0.35,15,MUTED)
chart_data = CategoryChartData(); chart_data.categories = ['LIGHT','PWR','DN','GenAgents']; chart_data.add_series('first playable mean (s)', (0.410,0.833,7.662,9.760))
chart = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.75), Inches(1.65), Inches(7.35), Inches(4.6), chart_data).chart
chart.has_legend=False
chart.value_axis.maximum_scale=10.5; chart.value_axis.minimum_scale=0
chart.value_axis.tick_labels.font.size=Pt(11); chart.category_axis.tick_labels.font.size=Pt(13)
chart.plots[0].vary_by_categories = True
# callouts right
card(s,8.65,1.72,3.95,0.9,fill=CREAM); tx(s,'DN',8.9,1.88,0.9,0.28,18,GREEN,True); tx(s,'7.662s mean / 17.716s p95；完整多模态链路的 first-playable proxy。',9.55,1.83,2.65,0.42,12,MUTED)
card(s,8.65,2.9,3.95,0.9,fill=CREAM); tx(s,'LIGHT',8.9,3.06,1.0,0.28,18,BLUE,True); tx(s,'0.410s mean；权威交互 baseline，但当前英文 adapter 语义贴合有限。',9.65,3.01,2.45,0.42,12,MUTED)
card(s,8.65,4.08,3.95,0.9,fill=CREAM); tx(s,'PWR',8.9,4.24,0.9,0.28,18,ORANGE,True); tx(s,'0.833s mean；轻量文本速度参考，不能直接代表完整游戏系统。',9.55,4.19,2.6,0.42,12,MUTED)
card(s,8.65,5.26,3.95,0.9,fill=CREAM); tx(s,'GenAgents',8.9,5.42,1.35,0.28,18,RED,True); tx(s,'9.760s mean；状态连续性补充 baseline，非完整多模态系统。',10.12,5.37,2.05,0.42,12,MUTED)
tx(s,'主结论：已具备“DN vs 外部系统”的效率对比，但不能把所有 baseline 说成完全同构。',0.78,6.85,11.5,0.3,13,INK,True)

# Slide 4 next steps
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
tx(s,'给老师汇报时的结论与下一步',0.65,0.45,8.3,0.55,32,INK,True)
tx(s,'效率部分已经能汇报；完整质量-效率权衡等质量组结果出来后再合并。',0.67,1.02,8.5,0.35,17,MUTED)
# three columns with open typography + rules
cols=[('现在可以说', '• DN 默认链路效率已量化\n• 外部系统 playable-latency 已补齐\n• LIGHT/PWR/GenAgents/WG 角色已区分', GREEN),
      ('必须说明边界', '• 当前只负责效率，不下质量优劣结论\n• DN row 是 first-playable proxy\n• baseline 不是完全同构系统', ORANGE),
      ('下一步', '• 冻结 Table 1 与补充表\n• 修正中文稿与图表说明\n• 等质量组补 judge/human/视觉质量后合并', BLUE)]
for i,(head,body,col) in enumerate(cols):
    x=0.85+i*4.1
    rule(s,x,2.05,2.45,col,4)
    tx(s,head,x,2.25,2.8,0.35,24,col,True)
    tx(s,body,x,2.9,3.05,1.45,15,INK)
# bottom statement
shp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(5.55), Inches(12.0), Inches(0.04))
shp.fill.solid(); shp.fill.fore_color.rgb=LINE; shp.line.fill.background()
tx(s,'建议汇报口径',0.75,5.8,2.0,0.3,17,GREEN,True)
tx(s,'“效率实验已完成第一版证据包：DN 主链路、外部 baseline 可玩延迟对比、预生成/readwait/Council 机制结果均已落盘；本阶段不做完整参数矩阵，避免只看速度而忽略质量。”',0.75,6.18,11.5,0.58,17,INK,True)

prs.save(pptx_path)

# Preview generation with PIL, matching slide content approximately.
try:
    font_title = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 54)
    font_h1 = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 42)
    font_h2 = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 28)
    font_body = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 23)
    font_small = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 17)
except Exception:
    font_title=font_h1=font_h2=font_body=font_small=ImageFont.load_default()

PW,PH=1600,900
C_BG=(248,244,235); C_INK=(35,39,47); C_MUTED=(95,101,113); C_GREEN=(56,132,103); C_BLUE=(54,103,161); C_ORANGE=(218,124,54); C_RED=(178,73,70); C_LINE=(213,205,188); C_CREAM=(255,251,242)

def draw_wrap(d, text, xy, font, fill, maxw, line_gap=6):
    x,y=xy
    for para in text.split('\n'):
        line=''
        for ch in para:
            test=line+ch
            if d.textbbox((0,0), test, font=font)[2] > maxw and line:
                d.text((x,y), line, font=font, fill=fill); y += font.size + line_gap; line=ch
            else:
                line=test
        if line:
            d.text((x,y), line, font=font, fill=fill); y += font.size + line_gap
    return y

def save_preview(idx, draw_func):
    img=Image.new('RGB',(PW,PH),C_BG); d=ImageDraw.Draw(img); draw_func(d)
    path=preview_dir/f'slide_{idx}.png'; img.save(path); return path

previews=[]

def p1(d):
    d.rectangle([0,0,26,PH], fill=C_GREEN)
    d.text((90,85),'DN 效率实验阶段性汇报',font=font_h2,fill=C_GREEN)
    draw_wrap(d,'已从“DN 自身跑通”\n推进到“外部 baseline 可玩延迟对比”',(90,185),font_title,C_INK,1000)
    d.text((95,410),'当前重点：只汇报效率，不提前下生成质量结论；质量评估后续共同完成。',font=font_body,fill=C_MUTED)
    for x,t,c in [(95,'DN fullchain 20/20 成功',C_GREEN),(425,'LIGHT 8-item 已跑通',C_BLUE),(730,'PWR / GenAgents 已入表',C_ORANGE)]:
        d.rounded_rectangle([x,545,x+285,606],radius=26,fill=c)
        d.text((x+18,562),t,font=font_small,fill='white')
    vals=[('DN',7.662,C_GREEN),('PWR',0.833,C_ORANGE),('LIGHT',0.410,C_BLUE),('GenAgents',9.760,C_RED)]
    d.text((1160,405),'first playable mean',font=font_small,fill=C_MUTED)
    for i,(lab,val,c) in enumerate(vals):
        x=1165+i*85; h=max(20,int(val/10*250)); d.rectangle([x,690-h,x+48,690],fill=c); d.text((x-5,708),lab,font=font_small,fill=C_MUTED); d.text((x-5,660-h),f'{val:.1f}s' if val>=1 else f'{val:.2f}s',font=font_small,fill=C_INK)

def p2(d):
    d.text((80,55),'现在已经做出的效率实验包',font=font_h1,fill=C_INK)
    d.text((82,120),'从默认链路、外部 baseline、机制消融到写作索引，已经形成可汇报材料。',font=font_body,fill=C_MUTED)
    nodes=[('1','DN 默认链路','20/20 成功\n效率画像完成',C_GREEN),('2','统一协议','playable latency\nsubset v1/v2',C_BLUE),('3','外部 baseline','LIGHT / PWR /\nGenAgents / WG',C_ORANGE),('4','写作包','Table 1 + caption\n局限性说明',C_RED)]
    for i,(num,title,body,c) in enumerate(nodes):
        x=105+i*370; d.ellipse([x,245,x+80,325],fill=c); d.text((x+28,258),num,font=font_h2,fill='white')
        if i<3: d.line([x+95,285,x+315,285],fill=C_LINE,width=4)
        d.text((x,355),title,font=font_h2,fill=C_INK); draw_wrap(d,body,(x,405),font_small,C_MUTED,240)
    for x,head,body,c in [(90,'已完成','DN 主链路、readwait/预生成、Council 部分消融、四类 baseline 延迟结果',C_GREEN),(570,'可汇报','主表候选、补充表、caption、局限性与汇报段落已经落盘',C_BLUE),(1050,'暂不做','完整参数矩阵/质量优劣结论：等后续质量评估合并后再判断',C_ORANGE)]:
        d.rounded_rectangle([x,585,x+380,735],radius=18,fill=C_CREAM,outline=C_LINE,width=2); d.text((x+22,605),head,font=font_h2,fill=c); draw_wrap(d,body,(x+22,650),font_small,C_MUTED,320)

def p3(d):
    d.text((80,55),'主实验表已补上外部系统效率对比',font=font_h1,fill=C_INK)
    d.text((82,120),'指标：first_playable_time_mean_s。数值越低越快，但系统能力范围不同。',font=font_small,fill=C_MUTED)
    labels=['LIGHT','PWR','DN','GenAgents']; vals=[0.410,0.833,7.662,9.760]; cols=[C_BLUE,C_ORANGE,C_GREEN,C_RED]
    x0,y0=150,690; maxv=10
    d.line([x0,y0,x0+720,y0],fill=C_LINE,width=3)
    for i,(lab,val,c) in enumerate(zip(labels,vals,cols)):
        y=220+i*105; w=int(val/maxv*650); d.text((x0-80,y+10),lab,font=font_small,fill=C_INK); d.rectangle([x0,y,x0+w,y+42],fill=c); d.text((x0+w+15,y+8),f'{val:.3f}s' if val<1 else f'{val:.3f}s',font=font_small,fill=C_INK)
    callouts=[('DN','完整多模态链路 first-playable proxy',C_GREEN),('LIGHT','权威交互 baseline，语义贴合有限',C_BLUE),('PWR','轻量文本速度参考',C_ORANGE),('GenAgents','状态连续性补充 baseline',C_RED)]
    for i,(h,b,c) in enumerate(callouts):
        y=215+i*125; d.rounded_rectangle([1010,y,1510,y+82],radius=16,fill=C_CREAM,outline=C_LINE,width=2); d.text((1035,y+18),h,font=font_body,fill=c); d.text((1160,y+20),b,font=font_small,fill=C_MUTED)
    d.text((90,815),'主结论：已具备“DN vs 外部系统”的效率对比，但不能把所有 baseline 说成完全同构。',font=font_body,fill=C_INK)

def p4(d):
    d.text((80,55),'给老师汇报时的结论与下一步',font=font_h1,fill=C_INK)
    d.text((82,120),'效率部分已经能汇报；完整质量-效率权衡等质量组结果出来后再合并。',font=font_body,fill=C_MUTED)
    cols=[('现在可以说','• DN 默认链路效率已量化\n• 外部系统 playable-latency 已补齐\n• LIGHT/PWR/GenAgents/WG 角色已区分',C_GREEN),('必须说明边界','• 当前只负责效率，不下质量优劣结论\n• DN row 是 first-playable proxy\n• baseline 不是完全同构系统',C_ORANGE),('下一步','• 冻结 Table 1 与补充表\n• 修正中文稿与图表说明\n• 等质量组补 judge/human/视觉质量后合并',C_BLUE)]
    for i,(head,body,c) in enumerate(cols):
        x=100+i*490; d.line([x,245,x+290,245],fill=c,width=6); d.text((x,275),head,font=font_h2,fill=c); draw_wrap(d,body,(x,350),font_small,C_INK,390)
    d.line([80,665,1520,665],fill=C_LINE,width=2); d.text((90,695),'建议汇报口径',font=font_body,fill=C_GREEN)
    draw_wrap(d,'“效率实验已完成第一版证据包：DN 主链路、外部 baseline 可玩延迟对比、预生成/readwait/Council 机制结果均已落盘；本阶段不做完整参数矩阵，避免只看速度而忽略质量。”',(90,740),font_body,C_INK,1370)

for i,fn in enumerate([p1,p2,p3,p4],1): previews.append(save_preview(i,fn))

# Validate PPTX package and slide count
with zipfile.ZipFile(pptx_path, 'r') as z:
    slides = [n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
    assert len(slides)==4, len(slides)
print('PPTX', pptx_path)
print('PREVIEWS')
for p in previews: print(p)
print('slide_count', len(slides))
