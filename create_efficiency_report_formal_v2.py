from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pathlib import Path
from project_paths import path_in_project
from PIL import Image, ImageDraw, ImageFont
import zipfile, textwrap

OUT = path_in_project('experiments', 'paper_method_view', '0_overview', 'teacher_report_efficiency_2026-04-26')
OUT.mkdir(parents=True, exist_ok=True)
pptx_path = OUT / 'DN_efficiency_progress_report_formal_v2_2026-04-26.pptx'
preview_dir = OUT / 'formal_v2_previews'
preview_dir.mkdir(exist_ok=True)

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# Palette: formal academic, high contrast
NAVY = RGBColor(27, 45, 72)
BLUE = RGBColor(39, 95, 153)
SKY = RGBColor(219, 234, 248)
TEAL = RGBColor(42, 128, 118)
GREEN = RGBColor(61, 139, 92)
ORANGE = RGBColor(213, 128, 56)
RED = RGBColor(182, 72, 72)
INK = RGBColor(35, 39, 47)
MUTED = RGBColor(91, 99, 112)
LIGHT = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 221, 230)
GREY = RGBColor(122, 130, 143)

FONT = 'Microsoft YaHei'

# ---------- helpers ----------
def slide_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.fill.background()
    # top rule
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.11))
    top.fill.solid(); top.fill.fore_color.rgb = NAVY; top.line.fill.background()
    return bg

def add_title(slide, title, subtitle=None, section='效率实验阶段性进展'):
    tx(slide, section, 0.55, 0.34, 3.4, 0.25, 10.5, BLUE, True)
    tx(slide, title, 0.55, 0.64, 9.7, 0.42, 25, INK, True)
    if subtitle:
        tx(slide, subtitle, 0.56, 1.1, 10.9, 0.33, 12.5, MUTED)
    rule(slide, 0.55, 1.52, 12.15, LINE, 1.1)

def tx(slide, text, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02); tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def para(slide, lines, x, y, w, h, size=12.5, color=INK, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.02); tf.margin_right=Inches(0.02); tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT; p.font.size=Pt(size); p.font.color.rgb=color
        if bullet:
            p.text = '• ' + line
    return box

def rule(slide, x, y, w, color=LINE, weight=1):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln

def panel(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(0.8)
    return shp

def badge(slide, text, x, y, w, color, font_size=10.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shp.fill.solid(); shp.fill.fore_color.rgb=color; shp.line.fill.background()
    tf=shp.text_frame; tf.clear(); tf.margin_left=Inches(0.07); tf.margin_right=Inches(0.07); tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(font_size); r.font.bold=True; r.font.color.rgb=WHITE
    return shp

def metric(slide, label, value, note, x, y, w, color):
    tx(slide, value, x, y, w, 0.42, 24, color, True)
    tx(slide, label, x, y+0.47, w, 0.22, 10.5, INK, True)
    tx(slide, note, x, y+0.74, w, 0.35, 9.5, MUTED)

def add_table(slide, rows, x, y, w, h, col_widths=None, header_fill=NAVY, font_size=9.5, header_size=9.5):
    r_count=len(rows); c_count=len(rows[0])
    shape=slide.shapes.add_table(r_count,c_count,Inches(x),Inches(y),Inches(w),Inches(h))
    table=shape.table
    if col_widths:
        for i,cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r,row in enumerate(rows):
        for c,val in enumerate(row):
            cell=table.cell(r,c)
            cell.text=str(val)
            cell.margin_left=Inches(0.05); cell.margin_right=Inches(0.05); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            cell.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name=FONT; run.font.size=Pt(header_size if r==0 else font_size)
                    run.font.bold = (r==0)
                    run.font.color.rgb = WHITE if r==0 else INK
            if r==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=header_fill
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255,255,255) if r%2 else RGBColor(241,245,249)
    return shape

# ---------- Slide 1 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
tx(s, 'DN 项目效率实验阶段性进展', 0.62, 0.62, 7.7, 0.55, 28, NAVY, True)
tx(s, 'Efficiency Experiment Progress Report · 2026-04-26', 0.64, 1.12, 5.2, 0.28, 12, MUTED)
rule(s, 0.62, 1.55, 5.6, BLUE, 2)
tx(s, '本阶段聚焦“效率”证据：端到端耗时、可玩内容返回延迟、外部系统对照与关键机制消融。生成质量、人评与图文一致性将在后续质量评估阶段合并。', 0.64, 1.83, 11.8, 0.64, 15, INK)
panel(s,0.7,2.86,3.0,1.45,WHITE,LINE); metric(s,'DN 默认 fullchain','20/20','完整链路成功；可作为效率主行',0.95,3.12,2.5,GREEN)
panel(s,3.95,2.86,3.0,1.45,WHITE,LINE); metric(s,'外部 baseline','4 类','LIGHT / PWR / GenAgents / WorldGeneration',4.20,3.12,2.45,BLUE)
panel(s,7.2,2.86,3.0,1.45,WHITE,LINE); metric(s,'统一协议','已建立','playable-latency schema + subset v1/v2',7.45,3.12,2.45,ORANGE)
panel(s,10.45,2.86,2.25,1.45,WHITE,LINE); metric(s,'写作包','已落盘','主表、caption、局限性说明',10.67,3.12,1.8,RED)
# bottom message
panel(s,0.72,5.05,11.95,1.32,RGBColor(236,242,248),RGBColor(194,208,226))
tx(s,'当前结论',0.96,5.25,1.4,0.28,14,BLUE,True)
tx(s,'效率实验已经从“DN 自身验证”推进到“DN 与外部系统在统一 playable-latency 口径下的对比”。但本阶段只讨论效率，不提前得出生成质量优劣结论。',2.05,5.17,9.75,0.55,17,INK,True)
tx(s,'建议定位：效率实验第一版证据包，而非最终完整质量-效率综合实验。',2.05,5.82,8.5,0.25,11.5,MUTED)

# ---------- Slide 2 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'实验口径：统一 playable-latency protocol','将不同系统统一转换为“从触发到返回可继续游玩内容”的效率指标。')
# Flow diagram panels
steps=[('输入集合','efficiency_playable_subset_v1/v2\n固定主题与触发输入，避免随意挑样本'),('系统运行','DN 默认链路 + 外部 baseline adapter\n记录 first playable 与 next turn'),('汇总指标','mean / p95 / success_rate\ncompleteness / continuity 作为效率护栏'),('论文表格','核心主表 + WorldGeneration 补充表\n明确每个 baseline 的角色边界')]
for i,(h,b) in enumerate(steps):
    x=0.75+i*3.12
    panel(s,x,2.02,2.62,1.42,WHITE,LINE)
    badge(s,str(i+1),x+0.12,2.18,0.36,[TEAL,BLUE,ORANGE,RED][i])
    tx(s,h,x+0.62,2.12,1.75,0.25,14,INK,True)
    tx(s,b,x+0.18,2.58,2.25,0.55,10.5,MUTED)
    if i<3:
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+2.68), Inches(2.72), Inches(x+3.0), Inches(2.72))
        ln.line.color.rgb=GREY; ln.line.width=Pt(1.5); ln.line.end_arrowhead=True
# metric definitions table
rows=[['指标','含义','当前用途'],['first_playable_time_s','首次返回可继续游玩的内容所需时间','主效率指标'],['p95_latency_s','高延迟尾部表现','观察稳定性/长尾风险'],['next_turn_time_s','后续一轮动作后的可玩返回时间','补充交互连续性'],['success_rate','运行是否成功返回可用内容','效率比较的前提条件'],['completeness / continuity','输出完整性与连续性代理指标','仅作效率护栏，不替代质量评价']]
add_table(s,rows,0.8,4.05,7.2,2.3,[1.75,3.05,2.4],font_size=8.7,header_size=9)
# right note
panel(s,8.45,4.05,3.85,2.3,RGBColor(255,252,244),RGBColor(229,205,166))
tx(s,'口径边界',8.72,4.27,1.6,0.28,15,ORANGE,True)
para(s,['DN 行使用 fullchain generate_option latency 作为 first-playable proxy。','不同 baseline 的系统形态不完全同构，因此主表解释重点是效率位置和系统角色。','生成质量、语义贴合、人评和图像一致性不在本阶段下结论。'],8.72,4.75,3.25,1.15,10.3,MUTED,bullet=True)
tx(s,'相关文件：baseline_integration/schema/playable_latency_run_schema.md；adapters/playable_protocol.py',0.8,6.92,11.6,0.25,8.5,MUTED)

# ---------- Slide 3 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'DN 默认配置效率画像：主链路已具备可量化结果','该页回答：DN 自身在默认配置下是否跑通、各阶段耗时在哪里。')
# left metrics
panel(s,0.75,1.86,3.0,1.25,WHITE,LINE); metric(s,'Worldview 生成','20/20','mean 33.657s · p95 68.872s',1.02,2.05,2.4,BLUE)
panel(s,0.75,3.35,3.0,1.25,WHITE,LINE); metric(s,'Fullchain 完整链路','20/20','worldview / option / 主角图均有耗时记录',1.02,3.54,2.4,GREEN)
panel(s,0.75,4.84,3.0,1.25,WHITE,LINE); metric(s,'效率护栏','0 fallback','图片返回、选项数、prompt 污染均有代理检查',1.02,5.03,2.4,ORANGE)
# right table
rows=[['阶段','样本','mean(s)','median(s)','p95(s)','说明'],['worldview default','20','33.657','23.664','68.872','结构化世界观生成'],['fullchain worldview','20','12.458','9.003','30.515','完整链路中的世界观阶段'],['generate option','20','7.662','0.022','17.716','当前 first-playable proxy'],['main character','20','55.912','56.546','77.015','主角图完成耗时']]
add_table(s,rows,4.25,1.95,8.05,2.55,[1.75,0.65,0.85,0.9,0.78,3.12],font_size=8.1,header_size=8.3)
# guardrail mini table
rows2=[['有效性代理','结果'],['worldview_success_rate','1.0'],['first_scene_success_rate','1.0'],['image_return_rate','1.0'],['option_count_ge_2_rate','1.0'],['fallback_trigger_rate','0.0'],['scene_prompt_pollution_rate','0.0']]
add_table(s,rows2,4.25,4.92,3.55,1.52,[2.45,0.9],font_size=7.5,header_size=8.0,header_fill=TEAL)
panel(s,8.15,4.92,4.15,1.52,RGBColor(238,244,250),RGBColor(196,210,226))
tx(s,'可支撑的效率结论',8.42,5.08,2.0,0.25,14,BLUE,True)
para(s,['DN 默认配置已有完整效率画像；','first-playable proxy 可用于外部系统效率对比；','质量判断仍需后续人评 / judge / 图文一致性评估。'],8.42,5.45,3.45,0.72,9.5,MUTED,bullet=True)

# ---------- Slide 4 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'外部 baseline 已补齐：每个系统有明确角色与证据边界','该页回答：现在已经和哪些外部论文/系统做了效率对照，以及它们各自能说明什么。')
rows=[['Baseline','当前角色','样本','已完成证据','主要限制'],
      ['LIGHT','权威互动对话 / game-world baseline','8','官方 checkpoint；8-item playable latency batch','英文 adapter；语义贴合 DN 中文主题有限'],
      ['Plan-Write-Revise','on-demand story generation 速度参考','5','官方模型包；5-item batch；summary 已生成','轻量文本系统，非完整多模态游戏链路'],
      ['GenAgents','状态连续性 / 多轮 agent 补充 baseline','8','8-item live run 转换为 playable-latency summary','不输出 DN-style 图像、世界观 JSON 或分支系统'],
      ['WorldGeneration','世界构造补充 baseline','8','fallback graph-to-playable 路径；8-item summary','非完整原论文 pipeline，建议放补充表'],
      ['AIDungeon / StoryDiffusion','当前周期排除或延期','-','status / decision note 已记录','legacy runtime 或 CUDA 硬件阻塞']]
add_table(s,rows,0.62,1.85,12.1,3.62,[1.75,2.2,0.55,3.1,4.5],font_size=7.6,header_size=8.0)
# role lane
for i,(name,col,txtt) in enumerate([('核心主表',BLUE,'DN + LIGHT + PWR + GenAgents'),('补充说明',ORANGE,'WorldGeneration fallback row'),('不进入当前效率主表',RED,'AIDungeon / StoryDiffusion')]):
    x=0.85+i*4.1
    panel(s,x,5.92,3.45,0.72,WHITE,LINE)
    badge(s,name,x+0.15,6.12,1.3,col,9)
    tx(s,txtt,x+1.58,6.1,1.7,0.28,9.5,MUTED,False)
tx(s,'整理原则：只有具备 source_links + protocol + raw_runs + summaries + status 的 baseline，才进入主表或补充表讨论。',0.75,6.9,11.7,0.25,9.5,MUTED)

# ---------- Slide 5 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'主表候选结果：DN 与外部系统的 first-playable 效率位置','主表可说明效率差异，但不能直接推导生成质量优劣。')
# chart
chart_data=CategoryChartData(); chart_data.categories=['LIGHT','PWR','DN','GenAgents']; chart_data.add_series('first_playable_time_mean_s',(0.410,0.833,7.662,9.760))
chart_shape=s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(0.72),Inches(1.9),Inches(6.0),Inches(3.35),chart_data)
chart=chart_shape.chart
chart.has_legend=False
chart.value_axis.maximum_scale=10.5; chart.value_axis.minimum_scale=0
chart.value_axis.tick_labels.font.size=Pt(8); chart.category_axis.tick_labels.font.size=Pt(9)
chart.plots[0].vary_by_categories=True
# table
rows=[['System','sample','first mean','p95','next turn','success','解释'],['DN','20','7.662','17.716','-','1.0','完整链路；使用 generate_option 作为 proxy'],['LIGHT','8','0.410','0.477','0.431','1.0','权威外部互动 baseline，语义贴合有限'],['PWR','5','0.833','0.963','0.837','1.0','轻量文本即时生成速度参考'],['GenAgents','8','9.760','18.697','6.285','0.875','多轮状态连续性补充 baseline']]
add_table(s,rows,7.0,1.88,5.45,3.35,[1.05,0.52,0.72,0.6,0.75,0.58,1.23],font_size=6.9,header_size=7.2,header_fill=BLUE)
# interpretation bands
panel(s,0.75,5.72,3.75,0.83,RGBColor(236,246,243),RGBColor(194,221,214))
tx(s,'可以得出的效率结论',0.98,5.92,1.8,0.25,12.5,TEAL,True)
tx(s,'DN 已有完整系统链路下的可玩返回效率；外部系统多数更快，但系统职责更窄。',0.98,6.21,3.05,0.26,8.7,MUTED)
panel(s,4.82,5.72,3.75,0.83,RGBColor(255,248,238),RGBColor(230,205,168))
tx(s,'不能直接得出的结论',5.05,5.92,1.8,0.25,12.5,ORANGE,True)
tx(s,'不能说 DN 质量优于或劣于所有 baseline；质量侧需后续统一评估。',5.05,6.21,3.05,0.26,8.7,MUTED)
panel(s,8.9,5.72,3.55,0.83,RGBColor(244,247,251),RGBColor(203,213,225))
tx(s,'主表写法',9.12,5.92,1.3,0.25,12.5,BLUE,True)
tx(s,'强调 unified protocol、角色边界和 latency evidence，而不是“谁全面胜出”。',9.12,6.21,2.85,0.26,8.7,MUTED)

# ---------- Slide 6 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'内部机制效率实验：已有结果解释 DN 延迟结构','该页用于说明：除了外部系统对比，DN 内部哪些设计影响效率。')
# three mechanism blocks
mechanisms=[('预生成 / readwait','60s 阅读窗口主结果','real-scene hit：33.3% → 83.3%\nsecond-click median：0.017s → 0.013s\n结论：收益体现在下一次交互，而非当前请求无条件更快。',GREEN),
            ('Council 开关','已有 fullchain / worldview 消融','默认 fullchain：20/20 成功\nno-council fullchain：19/20 成功\nworldview mean：12.458s vs 16.547s（当前结果显示 Council 影响稳定性与延迟结构）',BLUE),
            ('Pregeneration 清洗对照','clean12 pregen on/off','pregen_off worldview mean：73.377s\npregen_on worldview mean：16.566s\n结论：对主链路延迟控制有实质贡献，但需要结合 readwait 口径解释。',ORANGE)]
for i,(head,sub,body,col) in enumerate(mechanisms):
    x=0.75+i*4.12
    panel(s,x,1.95,3.52,3.65,WHITE,LINE)
    rule(s,x+0.18,2.18,1.05,col,3)
    tx(s,head,x+0.18,2.35,2.9,0.32,15,col,True)
    tx(s,sub,x+0.18,2.75,3.0,0.25,10.5,INK,True)
    tx(s,body,x+0.18,3.13,3.02,1.65,10.2,MUTED)
# bottom takeaway
panel(s,0.8,6.05,11.85,0.75,RGBColor(235,242,249),RGBColor(194,208,226))
tx(s,'效率部分当前建议表述',1.05,6.22,2.2,0.25,12.5,BLUE,True)
tx(s,'DN 的效率实验已经覆盖默认链路、外部系统可玩延迟对比和关键机制消融；完整“效率-质量权衡”需要等待后续质量评估结果合并。',3.0,6.18,8.9,0.35,12,INK,True)

# ---------- Slide 7 ----------
s=prs.slides.add_slide(prs.slide_layouts[6]); slide_bg(s)
add_title(s,'当前可交付材料与后续工作','该页用于收束汇报：哪些材料已经形成 source of truth，哪些工作留给下一阶段。')
# left deliverables
panel(s,0.75,1.9,5.55,3.9,WHITE,LINE)
tx(s,'已形成的材料',1.02,2.15,2.0,0.3,16,TEAL,True)
para(s,[
'主表：main_playable_latency_scaffold_2026-04-26.csv',
'含 LIGHT / WorldGeneration 版本与补充表：main_playable_latency_with_light_2026-04-26.csv；supplementary_playable_latency_worldgeneration_2026-04-26.csv',
'协议与子集：playable_latency_run_schema.md；playable_protocol.py；efficiency_playable_subset_v1/v2.json',
'写作包：experiment_section_draft、caption pack、limitations、main_experiment_writing_bundle_index',
'各 baseline 均有 status、raw_runs、summaries 或 decision note 支撑。'
],1.02,2.62,4.95,2.55,9.5,MUTED,bullet=True)
# right next
panel(s,6.75,1.9,5.55,3.9,WHITE,LINE)
tx(s,'下一阶段建议',7.02,2.15,2.0,0.3,16,BLUE,True)
para(s,[
'冻结效率主表口径：DN / LIGHT / PWR / GenAgents 作为核心主表，WorldGeneration 放补充表。',
'修正文稿表达：避免“全面优于 baseline”等不严谨说法，统一使用“效率位置”“系统角色”“口径边界”。',
'等待质量组补充：文本质量、人评、图文一致性、视觉质量、semantic fit，再合并成效率-质量综合结论。',
'完整 DN 参数矩阵暂不作为当前效率主线：只看速度会误导，需要质量结果共同解释。'
],7.02,2.62,4.95,2.55,9.5,MUTED,bullet=True)
# bottom formal closing
panel(s,0.75,6.18,11.55,0.62,RGBColor(244,247,251),RGBColor(203,213,225))
tx(s,'阶段性结论：效率实验第一版证据包已经完成，现阶段重点是冻结表格口径与规范解释；质量评估完成后，再形成最终综合实验结论。',1.0,6.36,10.95,0.25,11.5,INK,True)

prs.save(pptx_path)

# Validate package
with zipfile.ZipFile(pptx_path, 'r') as z:
    slides=[n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
    assert len(slides)==7

# ---------- PIL previews: manually approximate, enough for layout QA ----------
try:
    F_TITLE=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 38)
    F_H1=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 30)
    F_H2=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 22)
    F_BODY=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 17)
    F_SMALL=ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 13)
except Exception:
    F_TITLE=F_H1=F_H2=F_BODY=F_SMALL=ImageFont.load_default()
PW,PH=1600,900
C_NAVY=(27,45,72); C_BLUE=(39,95,153); C_TEAL=(42,128,118); C_GREEN=(61,139,92); C_ORANGE=(213,128,56); C_RED=(182,72,72); C_INK=(35,39,47); C_MUTED=(91,99,112); C_LIGHT=(246,248,251); C_WHITE=(255,255,255); C_LINE=(214,221,230)

def wrap_text(d,text,x,y,font,fill,maxw,line_gap=5):
    for para_text in text.split('\n'):
        line=''
        for ch in para_text:
            if d.textbbox((0,0),line+ch,font=font)[2] <= maxw or not line:
                line += ch
            else:
                d.text((x,y),line,font=font,fill=fill); y += font.size + line_gap; line=ch
        if line:
            d.text((x,y),line,font=font,fill=fill); y += font.size + line_gap
    return y

def base(title, subtitle=None):
    img=Image.new('RGB',(PW,PH),C_LIGHT); d=ImageDraw.Draw(img); d.rectangle([0,0,PW,13],fill=C_NAVY)
    d.text((66,38),'效率实验阶段性进展',font=F_SMALL,fill=C_BLUE)
    d.text((66,76),title,font=F_H1,fill=C_INK)
    if subtitle: d.text((66,130),subtitle,font=F_SMALL,fill=C_MUTED)
    d.line([66,180,1530,180],fill=C_LINE,width=2)
    return img,d

def draw_table(d, rows, x,y,w,h, widths=None, header=C_NAVY, fs=F_SMALL):
    n=len(rows); m=len(rows[0]); rh=h/n
    if widths is None: widths=[w/m]*m
    for r,row in enumerate(rows):
        cx=x
        for c,val in enumerate(row):
            cw=widths[c]
            fill=header if r==0 else ((255,255,255) if r%2 else (241,245,249))
            d.rectangle([cx,y+r*rh,cx+cw,y+(r+1)*rh],fill=fill,outline=C_LINE)
            color=(255,255,255) if r==0 else C_INK
            wrap_text(d,str(val),cx+7,y+r*rh+7,fs,color,cw-14,2)
            cx+=cw

def preview1():
    img,d=Image.new('RGB',(PW,PH),C_LIGHT),None
    img=Image.new('RGB',(PW,PH),C_LIGHT); d=ImageDraw.Draw(img); d.rectangle([0,0,PW,13],fill=C_NAVY)
    d.text((75,72),'DN 项目效率实验阶段性进展',font=F_TITLE,fill=C_NAVY)
    d.text((76,130),'Efficiency Experiment Progress Report · 2026-04-26',font=F_BODY,fill=C_MUTED)
    d.line([76,180,750,180],fill=C_BLUE,width=4)
    wrap_text(d,'本阶段聚焦“效率”证据：端到端耗时、可玩内容返回延迟、外部系统对照与关键机制消融。生成质量、人评与图文一致性将在后续质量评估阶段合并。',76,215,F_BODY,C_INK,1380)
    metrics=[('DN 默认 fullchain','20/20','完整链路成功；可作为效率主行',C_GREEN),('外部 baseline','4 类','LIGHT / PWR / GenAgents / WorldGeneration',C_BLUE),('统一协议','已建立','playable-latency schema + subset v1/v2',C_ORANGE),('写作包','已落盘','主表、caption、局限性说明',C_RED)]
    for i,(lab,val,note,col) in enumerate(metrics):
        x=85+i*385; d.rectangle([x,340,x+340,510],fill=C_WHITE,outline=C_LINE,width=2); d.text((x+25,365),val,font=F_TITLE,fill=col); d.text((x+25,425),lab,font=F_BODY,fill=C_INK); wrap_text(d,note,x+25,455,F_SMALL,C_MUTED,290)
    d.rectangle([85,615,1510,760],fill=(236,242,248),outline=(194,208,226),width=2)
    d.text((115,640),'当前结论',font=F_BODY,fill=C_BLUE)
    wrap_text(d,'效率实验已经从“DN 自身验证”推进到“DN 与外部系统在统一 playable-latency 口径下的对比”。但本阶段只讨论效率，不提前得出生成质量优劣结论。',250,635,F_H2,C_INK,1100)
    return img

def preview2():
    img,d=base('实验口径：统一 playable-latency protocol','将不同系统统一转换为“从触发到返回可继续游玩内容”的效率指标。')
    steps=['输入集合\nefficiency subset v1/v2','系统运行\nDN + baseline adapters','汇总指标\nmean / p95 / success','论文表格\n核心主表 + 补充表']
    for i,st in enumerate(steps):
        x=90+i*370; d.rectangle([x,250,x+310,395],fill=C_WHITE,outline=C_LINE,width=2); d.ellipse([x+15,270,x+55,310],fill=[C_TEAL,C_BLUE,C_ORANGE,C_RED][i]); d.text((x+27,278),str(i+1),font=F_SMALL,fill=(255,255,255)); wrap_text(d,st,x+75,268,F_BODY,C_INK,210)
        if i<3: d.line([x+318,322,x+360,322],fill=C_MUTED,width=3)
    rows=[['指标','含义','当前用途'],['first_playable_time_s','首次返回可继续游玩的内容所需时间','主效率指标'],['p95_latency_s','高延迟尾部表现','观察稳定性/长尾风险'],['next_turn_time_s','后续一轮动作后的可玩返回时间','补充交互连续性'],['success_rate','运行是否成功返回可用内容','效率比较前提'],['completeness / continuity','输出完整性与连续性代理指标','效率护栏，不替代质量评价']]
    draw_table(d,rows,95,490,870,280,[210,360,300])
    d.rectangle([1015,490,1480,770],fill=(255,252,244),outline=(229,205,166),width=2); d.text((1045,520),'口径边界',font=F_BODY,fill=C_ORANGE); wrap_text(d,'• DN 行使用 generate_option latency 作为 first-playable proxy。\n• 不同 baseline 不是完全同构系统。\n• 质量、人评、语义贴合不在本阶段下结论。',1045,570,F_SMALL,C_MUTED,380)
    return img

def preview3():
    img,d=base('DN 默认配置效率画像：主链路已具备可量化结果','DN 自身在默认配置下是否跑通、各阶段耗时在哪里。')
    for i,(v,l,n,c) in enumerate([('20/20','Worldview 生成','mean 33.657s · p95 68.872s',C_BLUE),('20/20','Fullchain 完整链路','worldview / option / 主角图均有记录',C_GREEN),('0 fallback','效率护栏','图片返回、选项数、prompt 污染均有代理检查',C_ORANGE)]):
        y=230+i*175; d.rectangle([90,y,445,y+135],fill=C_WHITE,outline=C_LINE,width=2); d.text((120,y+18),v,font=F_H1,fill=c); d.text((120,y+67),l,font=F_BODY,fill=C_INK); wrap_text(d,n,120,y+98,F_SMALL,C_MUTED,285)
    rows=[['阶段','样本','mean','median','p95','说明'],['worldview default','20','33.657','23.664','68.872','结构化世界观生成'],['fullchain worldview','20','12.458','9.003','30.515','完整链路中的世界观阶段'],['generate option','20','7.662','0.022','17.716','当前 first-playable proxy'],['main character','20','55.912','56.546','77.015','主角图完成耗时']]
    draw_table(d,rows,510,235,960,290,[210,80,100,100,100,370])
    rows2=[['有效性代理','结果'],['worldview_success_rate','1.0'],['first_scene_success_rate','1.0'],['image_return_rate','1.0'],['option_count_ge_2_rate','1.0'],['fallback_trigger_rate','0.0'],['scene_prompt_pollution_rate','0.0']]
    draw_table(d,rows2,510,590,420,190,[300,120],header=C_TEAL)
    d.rectangle([970,590,1470,780],fill=(238,244,250),outline=(196,210,226),width=2); d.text((1000,620),'可支撑的效率结论',font=F_BODY,fill=C_BLUE); wrap_text(d,'• DN 默认配置已有完整效率画像。\n• first-playable proxy 可用于外部系统效率对比。\n• 质量判断仍需后续评估。',1000,665,F_SMALL,C_MUTED,420)
    return img

def preview4():
    img,d=base('外部 baseline 已补齐：每个系统有明确角色与证据边界','现在已经和哪些外部论文/系统做了效率对照，以及它们各自能说明什么。')
    rows=[['Baseline','当前角色','样本','已完成证据','主要限制'],['LIGHT','权威互动对话 / game-world baseline','8','官方 checkpoint；8-item batch','英文 adapter；语义贴合 DN 中文主题有限'],['Plan-Write-Revise','on-demand story generation 速度参考','5','官方模型包；5-item summary','轻量文本系统，非完整多模态链路'],['GenAgents','状态连续性 / 多轮 agent 补充','8','8-item live run 转换 summary','不输出 DN-style 图像、世界观 JSON'],['WorldGeneration','世界构造补充 baseline','8','fallback graph-to-playable summary','非完整原论文 pipeline'],['AIDungeon / StoryDiffusion','当前周期排除或延期','-','status / decision note','legacy runtime 或 CUDA 阻塞']]
    draw_table(d,rows,75,220,1450,430,[185,270,70,395,530],fs=F_SMALL)
    for i,(h,t,c) in enumerate([('核心主表','DN + LIGHT + PWR + GenAgents',C_BLUE),('补充说明','WorldGeneration fallback row',C_ORANGE),('当前不进入','AIDungeon / StoryDiffusion',C_RED)]):
        x=100+i*490; d.rectangle([x,705,x+390,780],fill=C_WHITE,outline=C_LINE,width=2); d.rounded_rectangle([x+15,725,x+135,755],radius=12,fill=c); d.text((x+28,731),h,font=F_SMALL,fill=(255,255,255)); d.text((x+155,728),t,font=F_SMALL,fill=C_MUTED)
    return img

def preview5():
    img,d=base('主表候选结果：DN 与外部系统的 first-playable 效率位置','主表可说明效率差异，但不能直接推导生成质量优劣。')
    labels=['LIGHT','PWR','DN','GenAgents']; vals=[0.410,0.833,7.662,9.760]; cols=[C_BLUE,C_ORANGE,C_GREEN,C_RED]
    for i,(lab,val,c) in enumerate(zip(labels,vals,cols)):
        y=250+i*80; d.text((100,y+10),lab,font=F_SMALL,fill=C_INK); w=int(val/10*620); d.rectangle([210,y,210+w,y+35],fill=c); d.text((220+w,y+7),f'{val:.3f}s',font=F_SMALL,fill=C_INK)
    rows=[['System','sample','first mean','p95','next','success','解释'],['DN','20','7.662','17.716','-','1.0','完整链路 proxy'],['LIGHT','8','0.410','0.477','0.431','1.0','权威外部互动 baseline'],['PWR','5','0.833','0.963','0.837','1.0','轻量文本速度参考'],['GenAgents','8','9.760','18.697','6.285','0.875','多轮状态连续性补充']]
    draw_table(d,rows,815,230,680,310,[90,65,95,70,70,70,220],fs=F_SMALL,header=C_BLUE)
    for x,h,b,c in [(90,'可以得出的效率结论','DN 已有完整系统链路下的可玩返回效率；外部系统多数更快，但系统职责更窄。',C_TEAL),(575,'不能直接得出的结论','不能说 DN 质量优于或劣于所有 baseline；质量侧需后续统一评估。',C_ORANGE),(1060,'主表写法','强调 unified protocol、角色边界和 latency evidence。',C_BLUE)]:
        d.rectangle([x,670,x+410,780],fill=(255,255,255),outline=C_LINE,width=2); d.text((x+20,692),h,font=F_BODY,fill=c); wrap_text(d,b,x+20,730,F_SMALL,C_MUTED,350)
    return img

def preview6():
    img,d=base('内部机制效率实验：已有结果解释 DN 延迟结构','除了外部系统对比，DN 内部哪些设计影响效率。')
    blocks=[('预生成 / readwait','60s 阅读窗口主结果','real-scene hit：33.3% → 83.3%\nsecond-click median：0.017s → 0.013s\n收益体现在下一次交互。',C_GREEN),('Council 开关','fullchain / worldview 消融','默认 fullchain：20/20 成功\nno-council：19/20 成功\n影响稳定性与延迟结构。',C_BLUE),('Pregeneration 清洗对照','clean12 pregen on/off','pregen_off worldview mean：73.377s\npregen_on worldview mean：16.566s\n对主链路延迟控制有贡献。',C_ORANGE)]
    for i,(h,sub,body,c) in enumerate(blocks):
        x=90+i*490; d.rectangle([x,235,x+420,665],fill=C_WHITE,outline=C_LINE,width=2); d.line([x+25,265,x+160,265],fill=c,width=5); d.text((x+25,295),h,font=F_H2,fill=c); d.text((x+25,345),sub,font=F_BODY,fill=C_INK); wrap_text(d,body,x+25,395,F_BODY,C_MUTED,350)
    d.rectangle([90,725,1510,800],fill=(235,242,249),outline=(194,208,226),width=2); wrap_text(d,'效率部分当前建议表述：DN 的效率实验已经覆盖默认链路、外部系统可玩延迟对比和关键机制消融；完整“效率-质量权衡”需要等待后续质量评估结果合并。',115,744,F_BODY,C_INK,1350)
    return img

def preview7():
    img,d=base('当前可交付材料与后续工作','收束本阶段产物：哪些材料已经形成 source of truth，哪些工作留给下一阶段。')
    d.rectangle([90,230,750,690],fill=C_WHITE,outline=C_LINE,width=2); d.text((120,260),'已形成的材料',font=F_H2,fill=C_TEAL); wrap_text(d,'• 主表：main_playable_latency_scaffold_2026-04-26.csv\n• 补充表：main_playable_latency_with_light；supplementary_worldgeneration\n• 协议与子集：schema、playable_protocol、subset v1/v2\n• 写作包：实验章节草稿、caption、limitations、bundle index\n• baseline 均有 status、raw_runs、summaries 或 decision note。',120,315,F_SMALL,C_MUTED,570)
    d.rectangle([850,230,1510,690],fill=C_WHITE,outline=C_LINE,width=2); d.text((880,260),'下一阶段建议',font=F_H2,fill=C_BLUE); wrap_text(d,'• 冻结效率主表口径：DN / LIGHT / PWR / GenAgents 核心主表，WorldGeneration 补充。\n• 修正文稿表达：避免“全面优于 baseline”等不严谨说法。\n• 等待质量组补充文本质量、人评、图文一致性、视觉质量。\n• 完整 DN 参数矩阵暂不作为当前效率主线。',880,315,F_SMALL,C_MUTED,570)
    d.rectangle([90,740,1510,805],fill=(244,247,251),outline=(203,213,225),width=2); wrap_text(d,'阶段性结论：效率实验第一版证据包已经完成，现阶段重点是冻结表格口径与规范解释；质量评估完成后，再形成最终综合实验结论。',120,757,F_BODY,C_INK,1320)
    return img

funcs=[preview1,preview2,preview3,preview4,preview5,preview6,preview7]
preview_paths=[]
for i,fn in enumerate(funcs,1):
    img=fn(); p=preview_dir/f'slide_{i}.png'; img.save(p); preview_paths.append(p)
# montage
thumbs=[Image.open(p).resize((360,203)) for p in preview_paths]
cols=2; rows=(len(thumbs)+1)//2
mont=Image.new('RGB',(cols*380,rows*223),(235,238,242))
for i,img in enumerate(thumbs):
    mont.paste(img,(10+(i%cols)*380,10+(i//cols)*223))
montage=preview_dir/'montage.png'; mont.save(montage)

print('PPTX', pptx_path)
print('SLIDES', len(slides))
print('PREVIEW_DIR', preview_dir)
print('MONTAGE', montage)
